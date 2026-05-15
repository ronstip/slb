"""Scolto branding — logo and footer for slides.

Canonical source for the brand mark is the frontend Logo component
(frontend/src/components/Logo.tsx). Update both when the brand changes.
"""

import logging

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from api.agent.tools.presentation.theme import TemplateTheme, set_run_font

logger = logging.getLogger(__name__)

BRAND_NAME = "Scolto"
BRAND_DOT_COLOR = RGBColor(0xD9, 0x77, 0x57)  # matches Logo.tsx BRAND_DOT_COLOR


def draw_scolto_logo(
    slide,
    cx: int,
    cy: int,
    radius: int,
    theme: TemplateTheme,
    on_accent_bg: bool = False,
) -> None:
    """Draw the Scolto mark at (cx, cy) with the given half-size in EMU.

    Mirrors the SVG in Logo.tsx: four corner brackets framing a solid
    orange dot. Brackets use the theme's foreground color (or white on an
    accent background); the centre dot is always the fixed brand orange.
    """
    bracket_color = theme.white if on_accent_bg else theme.fg

    # Geometry derived from the 64-unit viewBox in Logo.tsx. `radius` is
    # the half-width of the logo's bounding box, so 1 viewBox unit = radius/32.
    inset = radius // 8                              # 4 viewBox units
    arm_len = int(radius * 7 / 16)                   # 14 viewBox units
    stroke_t = max(radius // 16, 9144)               # 2 viewBox units, min 1pt
    dot_r = max(int(radius * 7 / 32), 30000)         # 7 viewBox units

    left, right = cx - radius, cx + radius
    top, bottom = cy - radius, cy + radius

    def _rect(x, y, w, h, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    # Four corner brackets: each is a horizontal arm + a vertical arm
    # meeting at the bracket's elbow, drawn so the elbow is on the *inside*
    # of the corner (matching the SVG paths in Logo.tsx).
    tl_x, tl_y = left + inset, top + inset
    tr_x, tr_y = right - inset, top + inset
    br_x, br_y = right - inset, bottom - inset
    bl_x, bl_y = left + inset, bottom - inset

    _rect(tl_x, tl_y, arm_len, stroke_t, bracket_color)
    _rect(tl_x, tl_y, stroke_t, arm_len, bracket_color)
    _rect(tr_x - arm_len, tr_y, arm_len, stroke_t, bracket_color)
    _rect(tr_x - stroke_t, tr_y, stroke_t, arm_len, bracket_color)
    _rect(br_x - arm_len, br_y - stroke_t, arm_len, stroke_t, bracket_color)
    _rect(br_x - stroke_t, br_y - arm_len, stroke_t, arm_len, bracket_color)
    _rect(bl_x, bl_y - stroke_t, arm_len, stroke_t, bracket_color)
    _rect(bl_x, bl_y - arm_len, stroke_t, arm_len, bracket_color)

    # Centre dot (always the brand orange — matches Logo.tsx)
    oval = slide.shapes.add_shape(9, cx - dot_r, cy - dot_r, dot_r * 2, dot_r * 2)
    oval.fill.solid()
    oval.fill.fore_color.rgb = BRAND_DOT_COLOR
    oval.line.fill.background()


def add_footer(
    slide,
    theme: TemplateTheme,
    slide_width: int,
    slide_height: int,
    on_accent_bg: bool = False,
) -> None:
    """Add 'Powered by Scolto' footer with logo to bottom-right corner."""
    logo_r = Inches(0.115)
    logo_cx = slide_width - Inches(0.18) - logo_r
    logo_cy = slide_height - Inches(0.22)

    draw_scolto_logo(slide, logo_cx, logo_cy, logo_r, theme, on_accent_bg=on_accent_bg)

    text_color = theme.white if on_accent_bg else theme.muted
    text_w = Inches(1.85)
    txBox = slide.shapes.add_textbox(
        logo_cx - logo_r - text_w - Inches(0.04),
        slide_height - Inches(0.35),
        text_w, Inches(0.3),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.RIGHT
    run = para.add_run()
    run.text = f"Powered by {BRAND_NAME}"
    set_run_font(run, 7.5, theme, color=text_color)
