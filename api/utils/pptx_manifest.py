"""Extract a structured manifest from a .pptx template.

The manifest describes available layouts, placeholder types/positions,
and theme metadata (fonts, colors). It's stored in Firestore alongside
the template reference and injected into the agent's context so it knows
what layouts and styles are available when designing a deck.
"""

import io
import logging
from typing import Any

from lxml import etree
from pptx import Presentation

logger = logging.getLogger(__name__)

# Placeholder type enum values -> human-readable names
_PH_TYPE_NAMES = {
    0: "BODY",
    1: "TITLE",
    2: "BODY",
    3: "CENTER_TITLE",
    4: "SUBTITLE",
    5: "DATE",
    6: "SLIDE_NUMBER",
    7: "OBJECT",
    10: "PICTURE",
    12: "SLIDE_IMAGE",
    13: "SLIDE_NUMBER",
    14: "HEADER",
    15: "FOOTER",
    16: "DATE",
    18: "PICTURE",
}

# Placeholder types that are content slots (not metadata like date/footer/slide number)
_CONTENT_PH_TYPES = {0, 1, 2, 3, 4, 7, 10, 18}

# Structural signatures -> canonical layout names
# Signature = frozenset of content placeholder types
_STRUCTURAL_SIGNATURES = {
    frozenset([(3, 1), (4, 1)]): "Title Slide",          # CENTER_TITLE + SUBTITLE
    frozenset([(1, 1), (7, 1)]): "Title and Content",    # TITLE + 1 OBJECT
    frozenset([(1, 1), (2, 1)]): "Section Header",       # TITLE + 1 BODY
    frozenset([(1, 1), (7, 2)]): "Two Content",          # TITLE + 2 OBJECT
    frozenset([(1, 1), (2, 2), (7, 2)]): "Comparison",   # TITLE + 2 BODY + 2 OBJECT
    frozenset([(1, 1)]): "Title Only",                    # TITLE only
}


def _extract_theme(prs: Presentation) -> dict[str, Any]:
    """Extract font scheme and color scheme from the template's theme XML."""
    theme: dict[str, Any] = {
        "major_font": "Calibri Light",
        "minor_font": "Calibri",
        "colors": {},
    }
    try:
        for rel in prs.part.rels.values():
            if "slideMaster" not in str(rel.reltype):
                continue
            master_part = rel.target_part
            for mrel in master_part.rels.values():
                if "theme" not in str(mrel.reltype):
                    continue
                theme_blob = mrel.target_part.blob
                theme_xml = etree.fromstring(theme_blob)
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

                # Fonts
                font_scheme = theme_xml.find(".//a:fontScheme", ns)
                if font_scheme is not None:
                    major = font_scheme.find(".//a:majorFont/a:latin", ns)
                    minor = font_scheme.find(".//a:minorFont/a:latin", ns)
                    if major is not None and major.get("typeface"):
                        theme["major_font"] = major.get("typeface")
                    if minor is not None and minor.get("typeface"):
                        theme["minor_font"] = minor.get("typeface")

                # Colors
                color_scheme = theme_xml.find(".//a:clrScheme", ns)
                if color_scheme is not None:
                    for child in color_scheme:
                        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                        if len(child) > 0:
                            val_el = child[0]
                            val = val_el.get("val") or val_el.get("lastClr")
                            if val and val not in ("windowText", "window"):
                                theme["colors"][tag] = f"#{val}"
                            elif val == "windowText":
                                theme["colors"][tag] = "#000000"
                            elif val == "window":
                                theme["colors"][tag] = "#FFFFFF"
                return theme
    except Exception as e:
        logger.warning("pptx_manifest: theme extraction failed: %s", e)
    return theme


def _slot_name(ph_type_int: int, counters: dict[str, int]) -> str:
    """Assign a human-friendly slot name based on placeholder type.

    For layouts with 2 content placeholders (OBJECT or BODY), the first
    gets 'left' and the second gets 'right'. For a single content
    placeholder, it gets 'body'.
    """
    ph_type = _PH_TYPE_NAMES.get(ph_type_int, "UNKNOWN")

    if ph_type in ("CENTER_TITLE", "TITLE"):
        return "title"
    elif ph_type == "SUBTITLE":
        return "subtitle"
    elif ph_type in ("BODY", "OBJECT"):
        # Track all content placeholders together so that 2-BODY layouts
        # (Google Slides) and 2-OBJECT layouts (PowerPoint) both produce
        # left/right slots.
        count = counters.get("content", 0)
        counters["content"] = count + 1
        if count == 0:
            return "body"  # single content area — will be renamed to "left" if a second appears
        else:
            return "right"
    elif ph_type == "PICTURE":
        return "picture"
    return "unknown"


def extract_manifest(pptx_bytes: bytes) -> dict[str, Any]:
    """Extract a manifest from a .pptx file.

    Args:
        pptx_bytes: Raw bytes of the .pptx file.

    Returns:
        Manifest dict with layouts, theme, and slide dimensions.
    """
    prs = Presentation(io.BytesIO(pptx_bytes))

    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        counters: dict[str, int] = {}
        slots: list[str] = []
        placeholders: list[dict[str, Any]] = []

        for ph in layout.placeholders:
            ph_type_int = ph.placeholder_format.type.__int__() if hasattr(ph.placeholder_format.type, "__int__") else int(ph.placeholder_format.type)
            ph_idx = ph.placeholder_format.idx

            # Skip metadata placeholders (date, footer, slide number)
            if ph_type_int not in _CONTENT_PH_TYPES:
                continue

            slot = _slot_name(ph_type_int, counters)
            slots.append(slot)
            placeholders.append({
                "idx": ph_idx,
                "type": _PH_TYPE_NAMES.get(ph_type_int, "UNKNOWN"),
                "slot": slot,
                "left": ph.left,
                "top": ph.top,
                "width": ph.width,
                "height": ph.height,
            })

        # Fix: if we had 2 content placeholders, the first was named "body" but should be "left"
        if counters.get("content", 0) >= 2:
            for ph_info in placeholders:
                if ph_info["type"] in ("OBJECT", "BODY") and ph_info["slot"] == "body":
                    ph_info["slot"] = "left"
                    break  # only rename the first one
            slots = [ph_info["slot"] for ph_info in placeholders]

        layouts.append({
            "index": i,
            "name": layout.name,
            "slots": slots,
            "placeholders": placeholders,
        })

    theme = _extract_theme(prs)

    return {
        "layouts": layouts,
        "theme": theme,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
    }


def manifest_to_agent_context(manifest: dict, filename: str = "", gcs_path: str = "") -> str:
    """Format a manifest as a concise string for agent system prompt injection."""
    lines = []
    if filename:
        lines.append(f"Template: **{filename}**" + (f" (gcs_path: `{gcs_path}`)" if gcs_path else ""))
    lines.append("")
    lines.append("Layouts:")
    for layout in manifest.get("layouts", []):
        slots = layout.get("slots", [])
        slot_str = ", ".join(slots) if slots else "free canvas"
        name = layout["name"]
        extra = ""
        if not slots or all(s == "title" for s in slots):
            extra = " + free area for custom components"
        lines.append(f'- "{name}" [{slot_str}]{extra}')

    theme = manifest.get("theme", {})
    fonts = f"{theme.get('major_font', '?')}/{theme.get('minor_font', '?')}"
    colors = theme.get("colors", {})
    accent = colors.get("accent1", "?")
    lines.append(f"\nTheme: fonts={fonts}, accent={accent}")

    return "\n".join(lines)
