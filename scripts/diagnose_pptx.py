"""
Diagnostic script for the presentation rendering engine.

Tests each component type, layout matching, backward compatibility,
and roundtrip integrity. Output files go to tmp_pptx/ for visual inspection.

Usage:
    uv run python scripts/diagnose_pptx.py
"""

import io
import json
import sys
from pathlib import Path

# Setup path
project_root = str(Path(__file__).parent.parent)
sys.path.insert(0, project_root)

from pptx import Presentation

from api.agent.tools.presentation.branding import add_footer
from api.agent.tools.presentation.compat import convert_legacy_slides
from api.agent.tools.presentation.components import (
    fill_chart,
    fill_table,
    fill_text,
    render_key_finding,
    render_kpi_grid,
)
from api.agent.tools.presentation.manifest import (
    compute_free_area,
    find_blank_layout,
    get_slot_placeholders,
    resolve_layout,
)
from api.agent.tools.presentation.schemas import DeckPlan, parse_component
from api.agent.tools.presentation.theme import TemplateTheme
from api.agent.tools.presentation.validator import validate_deck_plan as _validate
from api.utils.pptx_manifest import extract_manifest

OUTPUT_DIR = Path(project_root) / "tmp_pptx"
TEMPLATE_PATH = Path(project_root) / "api" / "assets" / "templates" / "default.pptx"


def _load_template():
    """Load template, extract manifest, strip slides."""
    manifest = extract_manifest(TEMPLATE_PATH.read_bytes())
    prs = Presentation(str(TEMPLATE_PATH))
    xml_slides = prs.slides._sldIdLst
    rIds = [
        el.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        for el in list(xml_slides)
    ]
    for rId in rIds:
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
    for sld_id in list(xml_slides):
        xml_slides.remove(sld_id)
    return prs, manifest


def _save(prs, name: str) -> Path:
    buf = io.BytesIO()
    prs.save(buf)
    path = OUTPUT_DIR / name
    path.write_bytes(buf.getvalue())
    Presentation(io.BytesIO(buf.getvalue()))
    return path


def test_schema_validation():
    print("\n=== Schema Validation ===")
    plan = DeckPlan.model_validate({
        "title": "Test",
        "slides": [{"layout": "Title Slide", "content": {"title": {"component": "text", "text": "Hi"}}}],
    })
    print(f"  Valid plan: OK ({len(plan.slides)} slides)")

    try:
        parse_component({"component": "chart", "chart_type": "bar", "labels": ["A", "B"], "values": [1]})
        print("  ERROR: should have failed")
    except ValueError:
        print("  Mismatched chart data: caught OK")

    try:
        DeckPlan.model_validate({"title": "Test", "slides": []})
        print("  ERROR: should have failed")
    except ValueError:
        print("  Empty slides: caught OK")

    try:
        parse_component({"component": "kpi_grid", "items": [{"label": "x", "value": "y"}] * 9})
        print("  ERROR: should have failed")
    except ValueError:
        print("  KPI > 8: caught OK")


def test_layout_matching():
    print("\n=== Layout Matching ===")
    _, manifest = _load_template()
    layouts = manifest["layouts"]

    # Exact match (Anchor theme uses TITLE as layout name)
    idx = resolve_layout("TITLE", layouts)
    assert idx is not None, "TITLE layout not found"
    print(f"  Exact match (TITLE -> idx {idx}): OK")

    # Canonical name -> Google Slides alias
    idx = resolve_layout("Title and Content", layouts)
    assert idx is not None, "Title and Content not resolved"
    print(f"  Alias match (Title and Content -> idx {idx}): OK")

    # Two Content -> TITLE_AND_TWO_COLUMNS
    idx = resolve_layout("Two Content", layouts)
    assert idx is not None, "Two Content not resolved"
    print(f"  Alias match (Two Content -> idx {idx}): OK")

    idx = resolve_layout("Nonexistent Layout", layouts)
    assert idx is None, f"Expected None, got {idx}"
    print("  Unknown layout: None OK")

    idx = find_blank_layout(layouts)
    assert idx is not None, "Blank layout not found"
    print(f"  Blank fallback (idx {idx}): OK")


def test_individual_components():
    print("\n=== Individual Components ===")
    prs, manifest = _load_template()
    theme = TemplateTheme(manifest=manifest)
    layouts = manifest["layouts"]
    sw, sh = prs.slide_width, prs.slide_height

    # Title Slide
    idx = resolve_layout("Title Slide", layouts)
    slide = prs.slides.add_slide(prs.slide_layouts[idx])
    slot_phs = get_slot_placeholders(layouts[idx])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == slot_phs["title"]["idx"]:
            fill_text(ph, {"component": "text", "text": "Test Title", "style": "heading"}, theme)
        elif "subtitle" in slot_phs and ph.placeholder_format.idx == slot_phs["subtitle"]["idx"]:
            fill_text(ph, {"component": "text", "text": "Test Subtitle", "style": "subtitle"}, theme)
    add_footer(slide, theme, sw, sh)
    print("  Title Slide: OK")

    # Charts
    for chart_type in ("pie", "bar", "line"):
        idx = resolve_layout("Title and Content", layouts)
        slide = prs.slides.add_slide(prs.slide_layouts[idx])
        slot_phs = get_slot_placeholders(layouts[idx])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == slot_phs["title"]["idx"]:
                fill_text(ph, {"component": "text", "text": f"{chart_type.title()} Chart"}, theme)
            elif ph.placeholder_format.idx == slot_phs["body"]["idx"]:
                fill_chart(slide, ph, {"component": "chart", "chart_type": chart_type, "labels": ["A", "B", "C"], "values": [40, 35, 25]}, theme)
        add_footer(slide, theme, sw, sh)
        print(f"  Chart ({chart_type}): OK")

    # Table
    idx = resolve_layout("Title and Content", layouts)
    slide = prs.slides.add_slide(prs.slide_layouts[idx])
    slot_phs = get_slot_placeholders(layouts[idx])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == slot_phs["title"]["idx"]:
            fill_text(ph, {"component": "text", "text": "Data Table"}, theme)
        elif ph.placeholder_format.idx == slot_phs["body"]["idx"]:
            fill_table(slide, ph, {"component": "table", "columns": ["A", "B", "C"], "rows": [["1", "2", "3"], ["4", "5", "6"]]}, theme)
    add_footer(slide, theme, sw, sh)
    print("  Table: OK")

    # Two Content
    idx = resolve_layout("Two Content", layouts)
    slide = prs.slides.add_slide(prs.slide_layouts[idx])
    slot_phs = get_slot_placeholders(layouts[idx])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == slot_phs["title"]["idx"]:
            fill_text(ph, {"component": "text", "text": "Side-by-Side"}, theme)
        elif ph.placeholder_format.idx == slot_phs["left"]["idx"]:
            fill_chart(slide, ph, {"component": "chart", "chart_type": "bar", "labels": ["A", "B"], "values": [60, 40]}, theme)
        elif ph.placeholder_format.idx == slot_phs["right"]["idx"]:
            fill_chart(slide, ph, {"component": "chart", "chart_type": "pie", "labels": ["X", "Y"], "values": [70, 30]}, theme)
    add_footer(slide, theme, sw, sh)
    print("  Two Content: OK")

    # KPI Grid
    idx = resolve_layout("Title Only", layouts)
    slide = prs.slides.add_slide(prs.slide_layouts[idx])
    slot_phs = get_slot_placeholders(layouts[idx])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == slot_phs["title"]["idx"]:
            fill_text(ph, {"component": "text", "text": "Key Metrics"}, theme)
    free_area = compute_free_area(layouts[idx], sw, sh)
    render_kpi_grid(slide, {"component": "kpi_grid", "items": [
        {"label": "Posts", "value": "4,821"}, {"label": "Sentiment", "value": "0.62"},
        {"label": "Platform", "value": "TikTok"}, {"label": "Views", "value": "12.4M"},
    ]}, theme, free_area)
    add_footer(slide, theme, sw, sh)
    print("  KPI Grid: OK")

    # Key Finding
    idx = resolve_layout("Title Only", layouts)
    slide = prs.slides.add_slide(prs.slide_layouts[idx])
    slot_phs = get_slot_placeholders(layouts[idx])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == slot_phs["title"]["idx"]:
            fill_text(ph, {"component": "text", "text": "Key Finding"}, theme)
    free_area = compute_free_area(layouts[idx], sw, sh)
    render_key_finding(slide, {"component": "key_finding", "finding": "**52% negative** sentiment", "significance": "surprising"}, theme, free_area)
    add_footer(slide, theme, sw, sh)
    print("  Key Finding: OK")

    # Bullets
    idx = resolve_layout("Title and Content", layouts)
    slide = prs.slides.add_slide(prs.slide_layouts[idx])
    slot_phs = get_slot_placeholders(layouts[idx])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == slot_phs["title"]["idx"]:
            fill_text(ph, {"component": "text", "text": "Key Points"}, theme)
        elif ph.placeholder_format.idx == slot_phs["body"]["idx"]:
            fill_text(ph, {"component": "text", "bullets": ["**Bold point** with detail", "Second point", "Third point"]}, theme)
    add_footer(slide, theme, sw, sh)
    print("  Bullets: OK")

    path = _save(prs, "all_components.pptx")
    print(f"  Saved: {path} ({path.stat().st_size} bytes)")


def test_validator():
    print("\n=== Validator ===")

    result = _validate({
        "title": "Test", "collection_ids": ["col-123"],
        "slides": [
            {"layout": "Title Slide", "content": {"title": {"component": "text", "text": "Hello"}, "subtitle": {"component": "text", "text": "World"}}},
            {"layout": "Title and Content", "content": {"title": {"component": "text", "text": "Chart"}, "body": {"component": "chart", "chart_type": "pie", "labels": ["A", "B"], "values": [60, 40]}}},
        ],
    })
    print(f"  Valid plan: valid={result['valid']}, errors={len(result['errors'])}")
    assert result["valid"]

    result = _validate({
        "title": "Test",
        "slides": [{"layout": "Nonexistent Layout", "content": {"title": {"component": "text", "text": "X"}}}],
    })
    print(f"  Invalid layout: valid={result['valid']}, errors={len(result['errors'])}")
    assert not result["valid"]

    result = _validate({
        "title": "Test",
        "slides": [{"layout": "Section Header", "content": {"title": {"component": "text", "text": "S"}, "body": {"component": "chart", "chart_type": "bar", "labels": ["A"], "values": [10]}}}],
    })
    print(f"  Chart in BODY: warnings={len(result['warnings'])}")
    assert len(result["warnings"]) > 0

    print("  All validator tests: OK")


def test_backward_compat():
    print("\n=== Backward Compatibility ===")
    old_slides = [
        {"type": "title_slide", "title": "Legacy Title", "subtitle": "Sub"},
        {"type": "chart_pie", "title": "Pie", "labels": ["A", "B"], "values": [60, 40]},
        {"type": "bullets", "title": "Points", "bullets": ["One", "Two"]},
        {"type": "kpi_grid", "title": "KPIs", "items": [{"label": "X", "value": "1"}]},
        {"type": "key_finding", "title": "Find", "finding": "**Bold**", "significance": "notable"},
        {"type": "chart_row", "title": "Row", "charts": [
            {"type": "bar", "labels": ["A", "B"], "values": [10, 20]},
            {"type": "pie", "labels": ["X", "Y"], "values": [70, 30]},
        ]},
        {"type": "table", "title": "Tbl", "columns": ["A", "B"], "rows": [["1", "2"]]},
        {"type": "closing", "title": "Thanks", "message": "Done"},
    ]

    deck_plan = convert_legacy_slides(old_slides, "Legacy", ["col-123"])
    plan = DeckPlan.model_validate(deck_plan)
    print(f"  Converted: {len(old_slides)} -> {len(plan.slides)} slides")

    result = _validate(deck_plan)
    print(f"  Validation: valid={result['valid']}, errors={len(result['errors'])}")


def test_full_deck():
    print("\n=== Full Deck ===")
    deck_plan = {
        "title": "McDonald's Analysis",
        "collection_ids": ["col-abc"],
        "slides": [
            {"layout": "Title Slide", "content": {"title": {"component": "text", "text": "McDonald's TikTok Sentiment", "style": "heading"}, "subtitle": {"component": "text", "text": "Q1 2026 - 1,243 posts", "style": "subtitle"}}},
            {"layout": "Title Only", "content": {"title": {"component": "text", "text": "Key Metrics"}, "custom": {"component": "kpi_grid", "items": [{"label": "Posts", "value": "1,243"}, {"label": "Sentiment", "value": "+0.42"}, {"label": "Platform", "value": "TikTok"}, {"label": "Views", "value": "8.2M"}]}}},
            {"layout": "Title and Content", "content": {"title": {"component": "text", "text": "Sentiment"}, "body": {"component": "chart", "chart_type": "pie", "labels": ["Positive", "Neutral", "Negative"], "values": [523, 412, 308]}}},
            {"layout": "Two Content", "content": {"title": {"component": "text", "text": "Volume vs Sentiment"}, "left": {"component": "chart", "chart_type": "bar", "labels": ["TikTok", "Instagram", "Twitter"], "values": [680, 340, 143]}, "right": {"component": "chart", "chart_type": "pie", "labels": ["Pos", "Neg", "Neu"], "values": [42, 31, 27]}}},
            {"layout": "Title and Content", "content": {"title": {"component": "text", "text": "Top Channels"}, "body": {"component": "table", "columns": ["Channel", "Posts", "Sentiment"], "rows": [["@mcdonalds", "234", "Mixed"], ["@foodreviews", "189", "Negative"]]}}},
            {"layout": "Title Only", "content": {"title": {"component": "text", "text": "Key Finding"}, "custom": {"component": "key_finding", "finding": "**March 12 spike** correlates with price increase.", "significance": "surprising"}}},
            {"layout": "Title and Content", "content": {"title": {"component": "text", "text": "Recommendations"}, "body": {"component": "text", "bullets": ["**Monitor TikTok** closely", "Engage top creators", "Proactive messaging before price changes"]}}},
            {"layout": "Title Slide", "content": {"title": {"component": "text", "text": "Thank You", "style": "heading"}, "subtitle": {"component": "text", "text": "Analysis by Veille", "style": "subtitle"}}},
        ],
    }

    result = _validate(deck_plan)
    print(f"  Validation: valid={result['valid']}, errors={len(result['errors'])}, hints={len(result['optimization_hints'])}")

    prs, manifest = _load_template()
    theme = TemplateTheme(manifest=manifest)
    layouts = manifest["layouts"]
    sw, sh = prs.slide_width, prs.slide_height
    plan = DeckPlan.model_validate(deck_plan)

    for slide_spec in plan.slides:
        idx = resolve_layout(slide_spec.layout, layouts)
        if idx is None:
            idx = find_blank_layout(layouts)
        slide = prs.slides.add_slide(prs.slide_layouts[idx])
        layout_info = layouts[idx]
        slot_phs = get_slot_placeholders(layout_info)

        for slot_name, comp_spec in slide_spec.content.items():
            comp_type = comp_spec.get("component", "")
            if slot_name == "custom" or comp_type in ("kpi_grid", "key_finding"):
                free_area = compute_free_area(layout_info, sw, sh)
                if comp_type == "kpi_grid":
                    render_kpi_grid(slide, comp_spec, theme, free_area)
                elif comp_type == "key_finding":
                    render_key_finding(slide, comp_spec, theme, free_area)
                continue
            ph = None
            ph_info = slot_phs.get(slot_name)
            if ph_info:
                for p in slide.placeholders:
                    if p.placeholder_format.idx == ph_info["idx"]:
                        ph = p
                        break
            if ph is None:
                continue
            if comp_type == "text":
                fill_text(ph, comp_spec, theme)
            elif comp_type == "chart":
                fill_chart(slide, ph, comp_spec, theme)
            elif comp_type == "table":
                fill_table(slide, ph, comp_spec, theme)
        add_footer(slide, theme, sw, sh)

    path = _save(prs, "full_deck.pptx")
    print(f"  Full deck: {len(prs.slides)} slides, {path.stat().st_size} bytes")
    print(f"  Saved: {path}")


def main():
    OUTPUT_DIR.mkdir(exist_ok=True)
    print("PPTX Diagnostic — Component-Based Renderer")
    print("=" * 50)
    test_schema_validation()
    test_layout_matching()
    test_individual_components()
    test_validator()
    test_backward_compat()
    test_full_deck()
    print("\n" + "=" * 50)
    print("All tests passed!")
    print(f"Output files in: {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
