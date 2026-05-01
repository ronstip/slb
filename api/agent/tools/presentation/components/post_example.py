"""Post Examples component — embed real post media on a slide, hyperlinked to the source."""

import io
import logging
from datetime import datetime, timezone
from typing import Optional

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from api.agent.tools.presentation.post_lookup import (
    download_post_image,
    pick_primary_image,
)
from api.agent.tools.presentation.theme import TemplateTheme, set_run_font

logger = logging.getLogger(__name__)


def _truncate(text: Optional[str], max_chars: int) -> str:
    if not text:
        return ""
    text = text.strip().replace("\n", " ").replace("\r", " ")
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


def _time_ago(iso_str: Optional[str]) -> str:
    if not iso_str:
        return ""
    try:
        s = iso_str.replace("Z", "+00:00")
        if "+" not in s and "T" in s:
            s = s + "+00:00"
        ts = datetime.fromisoformat(s)
        if ts.tzinfo is None:
            ts = ts.replace(tzinfo=timezone.utc)
        delta = (datetime.now(timezone.utc) - ts).total_seconds()
        if delta < 60:
            return "now"
        if delta < 3600:
            return f"{int(delta // 60)}m"
        if delta < 86400:
            return f"{int(delta // 3600)}h"
        if delta < 86400 * 30:
            return f"{int(delta // 86400)}d"
        if delta < 86400 * 365:
            return f"{int(delta // (86400 * 30))}mo"
        return f"{int(delta // (86400 * 365))}y"
    except Exception:
        return ""


def _set_hyperlink(shape, url: Optional[str]) -> None:
    if not url:
        return
    try:
        shape.click_action.hyperlink.address = url
    except Exception as e:
        logger.debug("post_example: hyperlink failed: %s", e)


def _caption_text(post: dict) -> str:
    handle = (post.get("channel_handle") or "").strip()
    if handle and not handle.startswith("@"):
        handle = "@" + handle
    posted = _time_ago(post.get("posted_at"))
    return "  ·  ".join(p for p in (handle, posted) if p)


def _draw_caption(slide, x, y, w, h, text, theme, *, size, link_url=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    set_run_font(run, size, theme, color=theme.muted)
    if link_url:
        try:
            run.hyperlink.address = link_url
        except Exception:
            pass
    return box


def _draw_image_card(slide, x, y, w, h, post, image_bytes, theme, layout_mode):
    """Card centered on the post's saved image, with a caption strip and hyperlink."""
    post_url = post.get("post_url") or None
    caption = _caption_text(post)
    cap_size = 9 if layout_mode != "grid_3" else 8
    cap_h = Inches(0.28) if layout_mode != "grid_3" else Inches(0.24)
    img_h = h - (cap_h if caption else Inches(0.0))

    pic = None
    if image_bytes:
        try:
            pic = slide.shapes.add_picture(io.BytesIO(image_bytes), x, y, width=w, height=img_h)
        except Exception as e:
            logger.debug("post_example: add_picture failed, falling back: %s", e)

    if pic is None:
        # Image missing/invalid → render a content-text card so the slide still works.
        _draw_text_card(slide, x, y, w, h, post, theme, layout_mode)
        return

    _set_hyperlink(pic, post_url)

    if caption:
        cap_box = _draw_caption(
            slide, x, y + img_h, w, cap_h, caption, theme,
            size=cap_size, link_url=post_url,
        )
        _set_hyperlink(cap_box, post_url)


def _draw_text_card(slide, x, y, w, h, post, theme, layout_mode):
    """Fallback card for posts with no usable image — content text + handle, hyperlinked."""
    post_url = post.get("post_url") or None
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.surface
    bg.line.color.rgb = theme.border
    bg.line.width = Pt(0.5)
    _set_hyperlink(bg, post_url)

    pad = Inches(0.2)
    cap_h = Inches(0.28) if layout_mode != "grid_3" else Inches(0.24)
    body_size = {"single": 14, "grid_2": 11, "grid_3": 9}[layout_mode]
    body_max = {"single": 320, "grid_2": 200, "grid_3": 110}[layout_mode]

    title = (post.get("title") or "").strip()
    content = (post.get("content") or "").strip()
    body_text = _truncate(
        (title + " — " + content).strip(" —") if title and content else (title or content),
        body_max,
    )

    body_box = slide.shapes.add_textbox(
        x + pad, y + pad, w - pad * 2, h - pad * 2 - cap_h,
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.margin_left = Inches(0.0)
    body_tf.margin_right = Inches(0.0)
    body_tf.margin_top = Inches(0.0)
    body_tf.margin_bottom = Inches(0.0)
    para = body_tf.paragraphs[0]
    run = para.add_run()
    run.text = body_text or "(no content)"
    set_run_font(run, body_size, theme, color=theme.fg)
    if post_url:
        try:
            run.hyperlink.address = post_url
        except Exception:
            pass

    caption = _caption_text(post)
    if caption:
        _draw_caption(
            slide, x + pad, y + h - cap_h - pad, w - pad * 2, cap_h, caption, theme,
            size=9 if layout_mode != "grid_3" else 8, link_url=post_url,
        )


def render_post_examples(
    slide,
    spec: dict,
    theme: TemplateTheme,
    free_area: tuple,
    post_cache: Optional[dict] = None,
    image_cache: Optional[dict] = None,
) -> None:
    """Render up to 3 post-example cards in the free area.

    Each card shows the post's saved media (from gcs_uri) full-bleed, with a
    small caption and a hyperlink to the original post. Posts without media
    fall back to a content-text card.
    """
    post_cache = post_cache or {}
    image_cache = image_cache if image_cache is not None else {}

    layout_mode = spec.get("layout", "single")
    refs = spec.get("posts", []) or []
    if not refs:
        return

    posts: list[dict] = []
    for ref in refs:
        post_id = ref.get("post_id") if isinstance(ref, dict) else None
        if not post_id:
            posts.append({})
            continue
        posts.append(post_cache.get(post_id) or {"post_id": post_id})

    area_left, area_top, area_width, area_height = free_area
    area_height = min(area_height, Inches(5.0))

    if layout_mode == "single":
        cards = [(area_left, area_top, area_width, area_height)]
    elif layout_mode == "grid_2":
        gap = Inches(0.2)
        card_w = (area_width - gap) // 2
        cards = [
            (area_left, area_top, card_w, area_height),
            (area_left + card_w + gap, area_top, card_w, area_height),
        ]
    else:  # grid_3
        gap = Inches(0.15)
        card_w = (area_width - gap * 2) // 3
        cards = [
            (area_left + (card_w + gap) * i, area_top, card_w, area_height)
            for i in range(3)
        ]

    for (cx, cy, cw, ch), post in zip(cards, posts):
        post_id = post.get("post_id") if post else None
        if post_id and post_id not in image_cache:
            primary = pick_primary_image(post.get("media_refs") or [])
            image_cache[post_id] = download_post_image(primary) if primary else None
        image_bytes = image_cache.get(post_id) if post_id else None

        if image_bytes:
            _draw_image_card(slide, cx, cy, cw, ch, post or {}, image_bytes, theme, layout_mode)
        else:
            _draw_text_card(slide, cx, cy, cw, ch, post or {}, theme, layout_mode)
