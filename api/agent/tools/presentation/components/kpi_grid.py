"""KPI Grid component — custom card shapes in the free area below title."""

import logging

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from api.agent.tools.presentation.theme import TemplateTheme, set_run_font

logger = logging.getLogger(__name__)


def render_kpi_grid(
    slide,
    spec: dict,
    theme: TemplateTheme,
    free_area: tuple[int, int, int, int],
) -> None:
    """Render KPI cards as styled rectangles in the free area.

    Args:
        slide: The pptx slide object.
        spec: KPI grid spec with items: [{label, value}].
        theme: The active template theme.
        free_area: (left, top, width, height) in EMU defining available space.
    """
    items = spec.get("items", [])
    if not items:
        return
    items = items[:8]

    area_left, area_top, area_width, area_height = free_area
    count = len(items)
    per_row = min(count, 4)
    rows = (count + per_row - 1) // per_row

    gap = Inches(0.15)
    card_w = (area_width - gap * (per_row - 1)) // per_row
    card_h = Inches(2.2) if rows == 1 else Inches(1.9)
    v_gap = Inches(0.15)

    # Determine card colors based on background lightness
    # Use a slightly lighter/darker shade of the background for cards
    card_bg = theme.surface
    card_border = theme.border
    value_color = theme.white if theme.is_dark else theme.accent
    label_color = theme.muted

    for idx, item in enumerate(items):
        row = idx // per_row
        col = idx % per_row
        x = area_left + (card_w + gap) * col
        y = area_top + (card_h + v_gap) * row

        # Card background with rounded feel
        card = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = card_bg
        card.line.color.rgb = card_border
        card.line.width = Pt(0.5)

        # Thick accent strip at top
        strip = slide.shapes.add_shape(1, x, y, card_w, Inches(0.05))
        strip.fill.solid()
        strip.fill.fore_color.rgb = theme.accent
        strip.line.fill.background()

        # Value — large, bold, centered
        value_box = slide.shapes.add_textbox(
            x + Inches(0.15), y + Inches(0.25),
            card_w - Inches(0.3), Inches(1.1),
        )
        value_tf = value_box.text_frame
        value_tf.word_wrap = True
        value_para = value_tf.paragraphs[0]
        value_para.alignment = PP_ALIGN.CENTER
        value_run = value_para.add_run()
        value_run.text = str(item.get("value", ""))
        set_run_font(value_run, 30, theme, bold=True, color=value_color)

        # Label — smaller, muted, centered
        label_box = slide.shapes.add_textbox(
            x + Inches(0.1), y + card_h - Inches(0.65),
            card_w - Inches(0.2), Inches(0.5),
        )
        label_tf = label_box.text_frame
        label_tf.word_wrap = True
        label_para = label_tf.paragraphs[0]
        label_para.alignment = PP_ALIGN.CENTER
        label_run = label_para.add_run()
        label_run.text = str(item.get("label", ""))
        set_run_font(label_run, 11, theme, color=label_color)
