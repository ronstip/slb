"""Table component — renders a table at placeholder bounds."""

import logging
from typing import Any

from pptx.util import Pt

from api.agent.tools.presentation.theme import TemplateTheme, set_run_font

logger = logging.getLogger(__name__)


def fill_table(
    slide,
    placeholder,
    spec: dict,
    theme: TemplateTheme,
) -> None:
    """Render a table at the placeholder's bounds, then remove the placeholder.

    Args:
        slide: The pptx slide object.
        placeholder: The OBJECT placeholder whose bounds define table position.
        spec: Table component spec with columns and rows.
        theme: The active template theme.
    """
    columns = spec.get("columns", [])
    rows = spec.get("rows", [])

    if not columns or not rows:
        logger.warning("fill_table: missing columns or rows, skipping")
        return

    num_cols = len(columns)
    max_rows = min(len(rows), 12)
    num_rows = max_rows + 1  # +1 for header

    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height

    # Add table at placeholder bounds
    tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    tbl = tbl_shape.table

    # Header row
    for ci, col_name in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.accent
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.text = str(col_name)
        set_run_font(run, 11, theme, bold=True, color=theme.white)

    # Data rows
    for ri, row_data in enumerate(rows[:max_rows]):
        bg = theme.bg if ri % 2 == 0 else theme.surface
        for ci in range(num_cols):
            cell = tbl.cell(ri + 1, ci)
            val = row_data[ci] if ci < len(row_data) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            para = cell.text_frame.paragraphs[0]
            for run in para.runs:
                run.text = ""
            run = para.runs[0] if para.runs else para.add_run()
            run.text = str(val) if val is not None else ""
            set_run_font(run, 10, theme, color=theme.fg)

    # Remove the original placeholder element
    try:
        ph_el = placeholder._element
        ph_el.getparent().remove(ph_el)
    except Exception as e:
        logger.debug("fill_table: could not remove placeholder element: %s", e)
