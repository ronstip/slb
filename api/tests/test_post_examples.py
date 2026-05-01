"""Unit tests for the post_examples slide component."""

import io
from unittest.mock import MagicMock, patch

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pydantic import ValidationError

from api.agent.tools.presentation.schemas import (
    PostExamplesComponent,
    parse_component,
)
from api.agent.tools.presentation import post_lookup
from api.agent.tools.presentation.components import post_example
from api.agent.tools.presentation.components.post_example import (
    _truncate,
    render_post_examples,
)
from api.agent.tools.presentation.theme import TemplateTheme


# ── Schema validation ────────────────────────────────────────────────────────

def test_schema_rejects_zero_posts():
    with pytest.raises(ValidationError):
        PostExamplesComponent(layout="single", posts=[])


def test_schema_rejects_too_many_posts():
    posts = [{"post_id": f"p{i}", "collection_id": "c"} for i in range(4)]
    with pytest.raises(ValidationError):
        PostExamplesComponent(layout="grid_3", posts=posts)


def test_schema_rejects_layout_count_mismatch():
    with pytest.raises(ValidationError):
        PostExamplesComponent(
            layout="single",
            posts=[
                {"post_id": "p1", "collection_id": "c"},
                {"post_id": "p2", "collection_id": "c"},
            ],
        )


def test_schema_accepts_valid_grid_2():
    comp = PostExamplesComponent(
        layout="grid_2",
        posts=[
            {"post_id": "p1", "collection_id": "c"},
            {"post_id": "p2", "collection_id": "c"},
        ],
    )
    assert comp.layout == "grid_2"
    assert len(comp.posts) == 2


def test_schema_rejects_unknown_layout():
    with pytest.raises(ValidationError):
        PostExamplesComponent(
            layout="grid_4",
            posts=[{"post_id": "p1", "collection_id": "c"}],
        )


def test_parse_component_dispatches_post_examples():
    raw = {
        "component": "post_examples",
        "layout": "single",
        "posts": [{"post_id": "p1", "collection_id": "c1"}],
    }
    comp = parse_component(raw)
    assert isinstance(comp, PostExamplesComponent)


# ── BQ fetch ─────────────────────────────────────────────────────────────────

def test_fetch_posts_by_ids_empty_input_skips_bq():
    with patch.object(post_lookup, "get_bq") as mock_bq:
        result = post_lookup.fetch_posts_by_ids([])
    assert result == []
    mock_bq.assert_not_called()


def test_fetch_posts_by_ids_preserves_input_order():
    rows = [
        {"post_id": "p2", "platform": "twitter", "media_refs": []},
        {"post_id": "p1", "platform": "instagram", "media_refs": []},
        {"post_id": "p3", "platform": "tiktok", "media_refs": []},
    ]
    fake_bq = MagicMock()
    fake_bq.query.return_value = rows
    refs = [
        {"post_id": "p1", "collection_id": "c"},
        {"post_id": "p2", "collection_id": "c"},
        {"post_id": "p3", "collection_id": "c"},
    ]
    with patch.object(post_lookup, "get_bq", return_value=fake_bq):
        result = post_lookup.fetch_posts_by_ids(refs)
    assert [r["post_id"] for r in result] == ["p1", "p2", "p3"]
    sql = fake_bq.query.call_args[0][0]
    assert "IN UNNEST(@post_ids)" in sql


def test_fetch_posts_by_ids_parses_string_media_refs():
    rows = [
        {"post_id": "p1", "platform": "ig", "media_refs": '[{"gcs_uri":"gs://b/k"}]'}
    ]
    fake_bq = MagicMock()
    fake_bq.query.return_value = rows
    with patch.object(post_lookup, "get_bq", return_value=fake_bq):
        result = post_lookup.fetch_posts_by_ids([{"post_id": "p1", "collection_id": "c"}])
    assert isinstance(result[0]["media_refs"], list)
    assert result[0]["media_refs"][0]["gcs_uri"] == "gs://b/k"


def test_fetch_posts_by_ids_drops_unknown_ids():
    fake_bq = MagicMock()
    fake_bq.query.return_value = [{"post_id": "p1", "media_refs": []}]
    with patch.object(post_lookup, "get_bq", return_value=fake_bq):
        result = post_lookup.fetch_posts_by_ids(
            [
                {"post_id": "p1", "collection_id": "c"},
                {"post_id": "missing", "collection_id": "c"},
            ]
        )
    assert [r["post_id"] for r in result] == ["p1"]


def test_fetch_posts_by_ids_returns_empty_on_bq_failure():
    fake_bq = MagicMock()
    fake_bq.query.side_effect = RuntimeError("boom")
    with patch.object(post_lookup, "get_bq", return_value=fake_bq):
        result = post_lookup.fetch_posts_by_ids(
            [{"post_id": "p1", "collection_id": "c"}]
        )
    assert result == []


# ── Image helpers ────────────────────────────────────────────────────────────

def test_pick_primary_image_skips_videos():
    refs = [
        {"media_type": "video", "gcs_uri": "gs://b/v"},
        {"media_type": "image", "gcs_uri": "gs://b/i"},
    ]
    pick = post_lookup.pick_primary_image(refs)
    assert pick is not None
    assert pick["gcs_uri"] == "gs://b/i"


def test_pick_primary_image_returns_none_if_only_videos():
    refs = [{"media_type": "video", "gcs_uri": "gs://b/v"}]
    assert post_lookup.pick_primary_image(refs) is None


def test_pick_primary_image_accepts_default_media_type():
    refs = [{"gcs_uri": "gs://b/k"}]  # no media_type field
    pick = post_lookup.pick_primary_image(refs)
    assert pick is not None


def test_pick_primary_image_empty_input():
    assert post_lookup.pick_primary_image([]) is None
    assert post_lookup.pick_primary_image(None) is None


def test_transcode_webp_to_png():
    """Pillow can read WebP and we must re-encode to a python-pptx-supported format."""
    img = Image.new("RGB", (32, 32), (200, 100, 50))
    buf = io.BytesIO()
    img.save(buf, format="WEBP")
    webp_bytes = buf.getvalue()
    out = post_lookup._transcode_to_png(webp_bytes)
    assert out is not None
    # PNG magic header
    assert out[:4] == b"\x89PNG"


def test_transcode_drops_alpha_to_rgb():
    img = Image.new("RGBA", (16, 16), (0, 0, 0, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    out = post_lookup._transcode_to_png(buf.getvalue())
    assert out is not None
    reread = Image.open(io.BytesIO(out))
    assert reread.mode == "RGB"


def test_transcode_returns_none_on_garbage():
    assert post_lookup._transcode_to_png(b"not an image") is None


def test_is_natively_supported():
    img = Image.new("RGB", (8, 8), (1, 2, 3))
    buf = io.BytesIO()
    img.save(buf, format="JPEG")
    assert post_lookup._is_natively_supported(buf.getvalue()) is True

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    assert post_lookup._is_natively_supported(buf.getvalue()) is True

    buf = io.BytesIO()
    img.save(buf, format="WEBP")
    assert post_lookup._is_natively_supported(buf.getvalue()) is False


# ── Renderer helpers ─────────────────────────────────────────────────────────

def test_truncate_short_string_passthrough():
    assert _truncate("hello", 10) == "hello"


def test_truncate_long_string_appends_ellipsis():
    out = _truncate("a" * 50, 10)
    assert out.endswith("…")
    assert len(out) == 10


def test_truncate_handles_none():
    assert _truncate(None, 10) == ""


# ── Renderer integration (real Presentation, mocked image fetch) ─────────────

def _png_bytes() -> bytes:
    img = Image.new("RGB", (64, 64), (180, 180, 200))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_slide():
    prs = Presentation()
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


def test_image_post_sets_hyperlink_on_picture():
    slide = _make_slide()
    theme = TemplateTheme()
    post = {
        "post_id": "p1",
        "platform": "twitter",
        "channel_handle": "alice",
        "post_url": "https://twitter.com/alice/status/123",
        "posted_at": "2026-04-29T12:00:00Z",
        "media_refs": [{"gcs_uri": "gs://b/k", "media_type": "image"}],
    }
    spec = {"layout": "single", "posts": [{"post_id": "p1", "collection_id": "c"}]}
    free_area = (Inches(0.5), Inches(1.0), Inches(9.0), Inches(4.0))

    with patch.object(post_example, "download_post_image", return_value=_png_bytes()):
        render_post_examples(
            slide, spec, theme, free_area,
            post_cache={"p1": post}, image_cache={},
        )

    pictures = [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
    assert len(pictures) == 1
    assert pictures[0].click_action.hyperlink.address == "https://twitter.com/alice/status/123"


def test_text_only_post_falls_back_to_text_card_with_hyperlink():
    slide = _make_slide()
    theme = TemplateTheme()
    post = {
        "post_id": "p2",
        "platform": "twitter",
        "channel_handle": "bob",
        "content": "Just a plain text tweet, no media at all.",
        "post_url": "https://twitter.com/bob/status/456",
        "posted_at": "2026-04-30T08:00:00Z",
        "media_refs": [],
    }
    spec = {"layout": "single", "posts": [{"post_id": "p2", "collection_id": "c"}]}
    free_area = (Inches(0.5), Inches(1.0), Inches(9.0), Inches(4.0))

    render_post_examples(
        slide, spec, theme, free_area,
        post_cache={"p2": post}, image_cache={},
    )

    pictures = [s for s in slide.shapes if s.shape_type == 13]
    assert pictures == []
    # Find the rounded rectangle card and check it has the hyperlink.
    bg_shapes = [s for s in slide.shapes if s.has_text_frame is False or s.shape_type == 1]
    linked = [s for s in slide.shapes if s.click_action.hyperlink.address == post["post_url"]]
    assert linked, "expected at least one shape hyperlinked to the post URL"


def test_grid_3_renders_three_cards():
    slide = _make_slide()
    theme = TemplateTheme()
    posts = {
        f"p{i}": {
            "post_id": f"p{i}",
            "channel_handle": f"u{i}",
            "post_url": f"https://x.com/u{i}/status/{i}",
            "posted_at": "2026-04-30T08:00:00Z",
            "media_refs": [{"gcs_uri": f"gs://b/{i}", "media_type": "image"}],
        }
        for i in range(1, 4)
    }
    spec = {
        "layout": "grid_3",
        "posts": [{"post_id": f"p{i}", "collection_id": "c"} for i in range(1, 4)],
    }
    free_area = (Inches(0.5), Inches(1.0), Inches(9.0), Inches(4.0))

    with patch.object(post_example, "download_post_image", return_value=_png_bytes()):
        render_post_examples(
            slide, spec, theme, free_area,
            post_cache=posts, image_cache={},
        )

    pictures = [s for s in slide.shapes if s.shape_type == 13]
    assert len(pictures) == 3
    urls = {p.click_action.hyperlink.address for p in pictures}
    assert urls == {posts[f"p{i}"]["post_url"] for i in range(1, 4)}


def test_image_download_failure_falls_back_to_text_card():
    slide = _make_slide()
    theme = TemplateTheme()
    post = {
        "post_id": "p1",
        "channel_handle": "alice",
        "content": "Has media ref but download will fail.",
        "post_url": "https://example.com/post/1",
        "media_refs": [{"gcs_uri": "gs://b/k", "media_type": "image"}],
    }
    spec = {"layout": "single", "posts": [{"post_id": "p1", "collection_id": "c"}]}
    free_area = (Inches(0.5), Inches(1.0), Inches(9.0), Inches(4.0))

    with patch.object(post_example, "download_post_image", return_value=None):
        render_post_examples(
            slide, spec, theme, free_area,
            post_cache={"p1": post}, image_cache={},
        )

    pictures = [s for s in slide.shapes if s.shape_type == 13]
    assert pictures == []
    linked = [s for s in slide.shapes if s.click_action.hyperlink.address == post["post_url"]]
    assert linked
