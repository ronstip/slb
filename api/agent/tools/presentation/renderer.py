"""Render a PowerPoint presentation from a validated deck plan.

Template-native: uses the template's layouts and placeholders instead of
drawing raw shapes. Falls back to Blank layout with custom shapes when
no suitable layout is found.
"""

import io
import logging
import uuid
from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches

from api.agent.tools._idempotency import action_key, check_or_register
from api.agent.tools.presentation.branding import add_footer
from api.agent.tools.presentation.components import (
    fill_chart,
    fill_table,
    fill_text,
    render_key_finding,
    render_kpi_grid,
    render_post_examples,
)
from api.agent.tools.presentation.post_lookup import fetch_posts_by_ids
from api.agent.tools.presentation.manifest import (
    compute_free_area,
    find_blank_layout,
    get_slot_placeholders,
    resolve_layout,
)
from api.agent.tools.presentation.schemas import DeckPlan
from api.agent.tools.presentation.theme import TemplateTheme
from api.deps import get_gcs
from config.settings import get_settings

logger = logging.getLogger(__name__)

# Component types that are rendered as custom shapes in free area
_CUSTOM_COMPONENTS = {"kpi_grid", "key_finding", "post_examples"}

# Component dispatch for placeholder fills
_PLACEHOLDER_FILLERS = {
    "text": lambda slide, ph, spec, theme: fill_text(ph, spec, theme),
    "chart": fill_chart,
    "table": fill_table,
}

# Component dispatch for custom (free area) fills
_CUSTOM_FILLERS = {
    "kpi_grid": render_kpi_grid,
    "key_finding": render_key_finding,
    "post_examples": render_post_examples,
}


def _load_manifest(template_gcs_path: str, tool_context) -> dict:
    """Load manifest from session state or extract from template."""
    if tool_context is not None:
        state = tool_context.state
        ppt_template = state.get("ppt_template")
        if ppt_template and ppt_template.get("manifest"):
            return ppt_template["manifest"]

    from api.utils.pptx_manifest import extract_manifest

    if template_gcs_path:
        pptx_bytes = _download_template(template_gcs_path)
        return extract_manifest(pptx_bytes)

    template_path = _default_template_path()
    if template_path.exists():
        return extract_manifest(template_path.read_bytes())

    return {"layouts": [], "theme": {}, "slide_width": 12192000, "slide_height": 6858000}


def _default_template_path() -> Path:
    return Path(__file__).parent.parent.parent.parent / "assets" / "templates" / "default.pptx"


def _download_template(gcs_path: str) -> bytes:
    client = get_gcs()
    settings = get_settings()
    if gcs_path.startswith("gs://"):
        without_scheme = gcs_path[5:]
        bucket_name, _, blob_name = without_scheme.partition("/")
    else:
        bucket_name = settings.gcs_presentations_bucket or settings.gcs_exports_bucket
        blob_name = gcs_path
    return client.bucket(bucket_name).blob(blob_name).download_as_bytes()


def _get_bucket_name() -> str:
    settings = get_settings()
    return settings.gcs_presentations_bucket or settings.gcs_exports_bucket


def _upload_to_gcs(blob_name: str, data: bytes) -> str:
    client = get_gcs()
    bucket = client.bucket(_get_bucket_name())
    blob = bucket.blob(blob_name)
    blob.upload_from_string(
        data,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    return f"gs://{_get_bucket_name()}/{blob_name}"


def _load_presentation(template_gcs_path: str) -> Presentation:
    """Load the template presentation, preserving masters/layouts/theme."""
    try:
        if template_gcs_path:
            template_bytes = _download_template(template_gcs_path)
            prs = Presentation(io.BytesIO(template_bytes))
        else:
            prs = Presentation(str(_default_template_path()))

        # Remove existing slides but keep masters/layouts/theme
        xml_slides = prs.slides._sldIdLst
        rIds = [
            el.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            for el in list(xml_slides)
        ]
        for rId in rIds:
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass
        for sld_id in list(xml_slides):
            xml_slides.remove(sld_id)

        # Do NOT override slide dimensions - respect template
        return prs

    except Exception as e:
        logger.warning("generate_presentation: failed to load template: %s", e)
        prs = Presentation()
        return prs


def _find_placeholder_by_slot(slide, slot_name: str, slot_phs: dict) -> Optional[Any]:
    """Find the actual placeholder object on the slide matching a slot name."""
    ph_info = slot_phs.get(slot_name)
    if ph_info is None:
        return None
    target_idx = ph_info["idx"]
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == target_idx:
            return ph
    return None


def generate_presentation(
    deck_plan: dict = None,
    collection_ids: list[str] = None,
    title: str = "",
    template_gcs_path: str = "",
    slides: list[dict] = None,
    tool_context=None,
) -> dict:
    """Render a PowerPoint presentation from a deck plan or legacy slide spec.

    WHEN TO USE: After validating with validate_deck_plan. Pass the deck plan
    to render a presentation using the template's native layouts and components.

    DECK PLAN SCHEMA:
    {
      "title": "Deck Title",
      "collection_ids": ["col-xxx"],
      "template_gcs_path": "",
      "slides": [
        {
          "layout": "Title Slide",
          "content": {
            "title": {"component": "text", "text": "My Title"},
            "subtitle": {"component": "text", "text": "Subtitle here"}
          }
        },
        {
          "layout": "Title and Content",
          "content": {
            "title": {"component": "text", "text": "Slide Title"},
            "body": {"component": "chart", "chart_type": "pie", "labels": [...], "values": [...]}
          }
        },
        {
          "layout": "Two Content",
          "content": {
            "title": {"component": "text", "text": "Comparison"},
            "left": {"component": "chart", "chart_type": "bar", "labels": [...], "values": [...]},
            "right": {"component": "chart", "chart_type": "pie", "labels": [...], "values": [...]}
          }
        },
        {
          "layout": "Title Only",
          "content": {
            "title": {"component": "text", "text": "Key Metrics"},
            "custom": {"component": "kpi_grid", "items": [{"label": "...", "value": "..."}]}
          }
        },
        {
          "layout": "Title Only",
          "content": {
            "title": {"component": "text", "text": "Top Reactions"},
            "custom": {
              "component": "post_examples",
              "layout": "grid_3",
              "posts": [
                {"post_id": "p1", "collection_id": "c1"},
                {"post_id": "p2", "collection_id": "c1"},
                {"post_id": "p3", "collection_id": "c1"}
              ]
            }
          }
        }
      ]
    }

    COMPONENT TYPES:
    - text: {component: "text", text: "...", bullets: ["..."], style: "heading|body|subtitle"}
    - chart: {component: "chart", chart_type: "bar|pie|line", labels: [...], values: [...]}
    - table: {component: "table", columns: [...], rows: [[...]]}
    - kpi_grid: {component: "kpi_grid", items: [{label, value}]} - custom slot only
    - key_finding: {component: "key_finding", finding: "...", significance: "surprising|notable"} - custom slot only
    - post_examples: {component: "post_examples", layout: "single|grid_2|grid_3",
        posts: [{post_id, collection_id}]} - custom slot only.
        Tool fetches full post content + image from BQ at render time.
        Use exactly 1 post for "single", 2 for "grid_2", 3 for "grid_3".
        Pick post_ids the user has cited or that came back from a search query.

    LAYOUT GUIDE:
    - "Title Slide" [title, subtitle] - opening/closing
    - "Title and Content" [title, body] - single chart, table, or text
    - "Two Content" [title, left, right] - two charts or chart + text
    - "Section Header" [title, body] - section divider
    - "Comparison" [title, body, left, body_2, right] - labeled comparison
    - "Title Only" [title] + custom - kpi_grid, key_finding, post_examples

    Args:
        deck_plan: Structured deck plan (preferred).
        collection_ids: Collections this presentation covers.
        title: Presentation title.
        template_gcs_path: Optional GCS path to user template.
        slides: Legacy slide spec (backward compat - use deck_plan instead).

    Returns:
        status, presentation_id, title, collection_ids, slide_count on success.
    """
    # Handle backward compatibility
    if deck_plan is None and slides:
        from api.agent.tools.presentation.compat import convert_legacy_slides
        deck_plan = convert_legacy_slides(slides, title, collection_ids or [], template_gcs_path)

    if deck_plan is None:
        return {"status": "error", "message": "No deck_plan or slides provided."}

    # Parse plan
    try:
        plan = DeckPlan.model_validate(deck_plan)
    except Exception as e:
        return {"status": "error", "message": f"Invalid deck plan: {e}"}

    # Idempotency: an identical deck rendered earlier this session reuses
    # its presentation_id instead of producing a duplicate file.
    _idempo_key = action_key("generate_presentation", {
        "deck_plan": plan.model_dump(exclude_none=True),
        "title": title or "",
        "template_gcs_path": template_gcs_path or "",
        "collection_ids": sorted(collection_ids or []),
    })
    _existing = check_or_register(tool_context, _idempo_key, dry_run=True)
    if _existing:
        return {
            "status": "duplicate",
            "presentation_id": _existing["artifact_id"],
            "message": (
                "An identical presentation was already rendered earlier in this session - "
                "reusing it. Don't render another."
            ),
        }

    effective_title = title or plan.title or "Presentation"
    effective_collections = collection_ids or plan.collection_ids or []
    effective_template = template_gcs_path or plan.template_gcs_path or ""

    # Build theme from manifest + session state
    try:
        _ctx = tool_context
    except Exception:
        _ctx = None
    manifest = _load_manifest(effective_template, _ctx)
    theme = TemplateTheme.from_session(_ctx, manifest=manifest)

    manifest_layouts = manifest.get("layouts", [])
    slide_width = manifest.get("slide_width", 12192000)
    slide_height = manifest.get("slide_height", 6858000)

    # Load presentation template
    prs = _load_presentation(effective_template)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Pre-fetch all post examples in a single BQ call so each post_id only
    # hits BigQuery once across the whole deck. The per-call image cache is
    # filled lazily inside the renderer.
    all_post_refs: list[dict] = []
    for slide_spec in plan.slides:
        for raw in slide_spec.content.values():
            if isinstance(raw, dict) and raw.get("component") == "post_examples":
                for ref in raw.get("posts", []) or []:
                    if isinstance(ref, dict) and ref.get("post_id"):
                        all_post_refs.append(ref)

    post_cache: dict = {}
    if all_post_refs:
        agent_id = None
        if tool_context is not None:
            try:
                agent_id = tool_context.state.get("active_agent_id")
            except Exception:
                agent_id = None
        for row in fetch_posts_by_ids(all_post_refs, agent_id=agent_id):
            pid = row.get("post_id")
            if pid:
                post_cache[pid] = row
    image_cache: dict = {}

    # Render slides
    rendered = 0
    for slide_spec in plan.slides:
        try:
            # Resolve layout
            layout_idx = resolve_layout(slide_spec.layout, manifest_layouts)
            if layout_idx is None:
                layout_idx = find_blank_layout(manifest_layouts)

            # Get layout from prs
            if layout_idx < len(prs.slide_layouts):
                layout = prs.slide_layouts[layout_idx]
            else:
                layout = prs.slide_layouts[-1]

            slide = prs.slides.add_slide(layout)

            # Get layout info from manifest for slot mapping
            layout_info = None
            for li in manifest_layouts:
                if li["index"] == layout_idx:
                    layout_info = li
                    break
            if layout_info is None:
                layout_info = {"placeholders": [], "slots": []}

            slot_phs = get_slot_placeholders(layout_info)

            # Fill each content slot
            for slot_name, comp_spec in slide_spec.content.items():
                comp_type = comp_spec.get("component", "")

                if slot_name == "custom" or comp_type in _CUSTOM_COMPONENTS:
                    # Render in free area
                    free_area = compute_free_area(layout_info, slide_width, slide_height)
                    filler = _CUSTOM_FILLERS.get(comp_type)
                    if filler is render_post_examples:
                        filler(
                            slide, comp_spec, theme, free_area,
                            post_cache=post_cache, image_cache=image_cache,
                        )
                    elif filler:
                        filler(slide, comp_spec, theme, free_area)
                    continue

                # Find the actual placeholder on the slide
                ph = _find_placeholder_by_slot(slide, slot_name, slot_phs)
                if ph is None:
                    logger.debug(
                        "render: slot %r not found on slide (layout=%r), skipping",
                        slot_name, slide_spec.layout,
                    )
                    continue

                # Dispatch to component filler
                # Charts and tables work with any placeholder type - they
                # render at the placeholder's bounds and remove its XML element.
                # This handles both OBJECT placeholders (python-pptx templates)
                # and BODY placeholders (Google Slides templates).
                if comp_type == "text":
                    fill_text(ph, comp_spec, theme)
                elif comp_type == "chart":
                    fill_chart(slide, ph, comp_spec, theme)
                elif comp_type == "table":
                    fill_table(slide, ph, comp_spec, theme)
                else:
                    logger.warning(
                        "render: unknown component %r in slot %r", comp_type, slot_name
                    )

            # Add branding footer
            add_footer(slide, theme, slide_width, slide_height)
            rendered += 1

        except Exception as e:
            logger.warning(
                "generate_presentation: failed to render slide %r: %s",
                slide_spec.layout, e,
            )

    if rendered == 0:
        return {"status": "error", "message": "No slides could be rendered - check the deck plan."}

    # Serialize and upload
    presentation_id = f"ppt-{uuid.uuid4().hex[:10]}"
    check_or_register(tool_context, _idempo_key, artifact_id=presentation_id)
    safe_title = (effective_title).replace(" ", "_").replace("/", "-")[:60]
    blob_name = f"presentations/{presentation_id}/{safe_title}.pptx"

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    try:
        _upload_to_gcs(blob_name, pptx_bytes)
        logger.info("generate_presentation: uploaded %s (%d bytes)", blob_name, len(pptx_bytes))
    except Exception as e:
        logger.error("generate_presentation: GCS upload failed: %s", e)
        return {"status": "error", "message": f"Failed to save presentation: {e}"}

    return {
        "status": "success",
        "presentation_id": presentation_id,
        "title": effective_title,
        "collection_ids": effective_collections,
        "slide_count": rendered,
        "gcs_path": blob_name,
        "message": f"Presentation created with {rendered} slides. The user can download it from the Artifacts panel.",
    }
