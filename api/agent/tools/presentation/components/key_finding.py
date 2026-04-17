"""Key Finding component — accent-bordered card with finding text."""

import logging

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from api.agent.tools.presentation.theme import (
    TemplateTheme,
    parse_md_segments,
    set_run_font,
)

logger = logging.getLogger(__name__)


def render_key_finding(
    slide,
    spec: dict,
    theme: TemplateTheme,
    free_area: tuple[int, int, int, int],
) -> None:
    """Render a key finding card in the free area.

    Args:
        slide: The pptx slide object.
        spec: Key finding spec with finding text and significance.
        theme: The active template theme.
        free_area: (left, top, width, height) in EMU defining available space.
    """
    finding = spec.get("finding", "")
    significance = spec.get("significance", "notable")

    area_left, area_top, area_width, area_height = free_area

    # Use most of the area but not all — leave breathing room
    card_margin = Inches(0.1)
    card_left = area_left + card_margin
    card_top = area_top + card_margin
    card_width = area_width - card_margin * 2
    card_height = min(area_height - card_margin * 2, Inches(3.5))

    # Card background
    card = slide.shapes.add_shape(1, card_left, card_top, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = theme.surface
    border_color = theme.accent if significance == "surprising" else theme.border
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5 if significance == "surprising" else 0.75)

    # Accent strip at top of card
    strip = slide.shapes.add_shape(1, card_left, card_top, card_width, Inches(0.05))
    strip.fill.solid()
    strip.fill.fore_color.rgb = theme.accent
    strip.line.fill.background()

    # Badge
    badge_text = "SURPRISING SIGNAL" if significance == "surprising" else "KEY FINDING"
    badge_box = slide.shapes.add_textbox(
        card_left + Inches(0.3), card_top + Inches(0.25),
        card_width - Inches(0.6), Inches(0.3),
    )
    badge_tf = badge_box.text_frame
    badge_para = badge_tf.paragraphs[0]
    badge_run = badge_para.add_run()
    badge_run.text = badge_text
    set_run_font(badge_run, 10, theme, bold=True, color=theme.accent)

    # Finding body with markdown support — use larger text for impact
    if finding:
        body_box = slide.shapes.add_textbox(
            card_left + Inches(0.3), card_top + Inches(0.7),
            card_width - Inches(0.6), card_height - Inches(1.0),
        )
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        para = body_tf.paragraphs[0]
        # Use fg color for text on the card surface
        text_color = theme.fg
        for segment, is_bold in parse_md_segments(finding):
            run = para.add_run()
            run.text = segment
            set_run_font(
                run, 18, theme,
                bold=is_bold,
                color=theme.accent if is_bold else text_color,
            )
