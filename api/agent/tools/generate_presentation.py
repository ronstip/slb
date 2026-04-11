"""
Generate Presentation Tool — agent-driven PowerPoint deck builder.

Renders a branded .pptx from the agent's structured `slides` spec.
Colors are derived from the user's selected accent color and light/dark theme
preference, both of which are passed through session state.
"""

import io
import logging
import math
import re
import uuid
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from api.deps import get_gcs
from config.settings import get_settings

logger = logging.getLogger(__name__)

# ── Default accent (matches DEFAULT_ACCENT in accent-colors.ts) ───────────────
_DEFAULT_ACCENT_HEX = "#4A7C8F"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
_FONT = "Inter"
_FONT_FALLBACK = "Segoe UI"


# ── Color helpers ─────────────────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "4A7C8F"
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _hex_to_rgb_color(hex_color: str) -> RGBColor:
    r, g, b = _hex_to_rgb(hex_color)
    return _rgb(r, g, b)


def _rgb_to_hsl(r: int, g: int, b: int) -> tuple[float, float, float]:
    rf, gf, bf = r / 255, g / 255, b / 255
    mx, mn = max(rf, gf, bf), min(rf, gf, bf)
    l = (mx + mn) / 2
    if mx == mn:
        return 0.0, 0.0, l
    d = mx - mn
    s = d / (2 - mx - mn) if l > 0.5 else d / (mx + mn)
    if mx == rf:
        h = ((gf - bf) / d + (6 if gf < bf else 0)) / 6
    elif mx == gf:
        h = ((bf - rf) / d + 2) / 6
    else:
        h = ((rf - gf) / d + 4) / 6
    return h * 360, s, l


def _hsl_to_rgb(h: float, s: float, l: float) -> tuple[int, int, int]:
    h = ((h % 360) + 360) % 360
    a = s * min(l, 1 - l)
    def f(n: float) -> int:
        k = (n + h / 30) % 12
        c = l - a * max(min(k - 3, 9 - k, 1), -1)
        return round(255 * max(0, min(1, c)))
    return f(0), f(8), f(4)


def _build_chart_palette(accent_hex: str, is_dark: bool) -> list[RGBColor]:
    """Generate 5-color monochromatic palette from accent — mirrors accent-colors.ts."""
    r, g, b = _hex_to_rgb(accent_hex)
    h, s, _ = _rgb_to_hsl(r, g, b)
    shades = (
        [(0.55, 1.0), (0.70, 0.75), (0.40, 1.1), (0.62, 0.60), (0.48, 0.85)]
        if is_dark else
        [(0.35, 0.90), (0.50, 0.75), (0.25, 1.0), (0.62, 0.55), (0.42, 0.65)]
    )
    palette = []
    for l_val, s_factor in shades:
        rs, gs, bs = _hsl_to_rgb(h, min(s * s_factor, 0.85), l_val)
        palette.append(_rgb(rs, gs, bs))
    return palette


class _Theme:
    """All slide colors derived from user accent + light/dark preference."""

    def __init__(self, accent_hex: str = _DEFAULT_ACCENT_HEX, is_dark: bool = False):
        self.accent = _hex_to_rgb_color(accent_hex)
        self.is_dark = is_dark

        if is_dark:
            self.bg      = _rgb(0x0A, 0x0A, 0x0A)
            self.card    = _rgb(0x17, 0x17, 0x17)
            self.surface = _rgb(0x1A, 0x1A, 0x1A)
            self.border  = _rgb(0x2E, 0x2E, 0x2E)
            self.fg      = _rgb(0xED, 0xED, 0xED)
            self.muted   = _rgb(0xA1, 0xA1, 0xA1)
            self.white   = _rgb(0xFF, 0xFF, 0xFF)
        else:
            self.bg      = _rgb(0xFA, 0xFA, 0xFA)
            self.card    = _rgb(0xFF, 0xFF, 0xFF)
            self.surface = _rgb(0xF0, 0xF0, 0xF0)
            self.border  = _rgb(0xD4, 0xD4, 0xD4)
            self.fg      = _rgb(0x0A, 0x0A, 0x0A)
            self.muted   = _rgb(0x52, 0x52, 0x52)
            self.white   = _rgb(0xFF, 0xFF, 0xFF)

        # Derive a darker shade of accent for bold text on light backgrounds
        r, g, b = _hex_to_rgb(accent_hex)
        h, s, l = _rgb_to_hsl(r, g, b)
        dr, dg, db = _hsl_to_rgb(h, min(s * 1.1, 1.0), max(l - 0.12, 0.15))
        self.accent_dark = _rgb(dr, dg, db)

        self.chart_palette = _build_chart_palette(accent_hex, is_dark)

    @classmethod
    def from_session(cls, tool_context) -> "_Theme":
        """Build theme from ADK session state (falls back to defaults)."""
        if tool_context is None:
            return cls()
        state = tool_context.state
        accent = state.get("accent_color") or _DEFAULT_ACCENT_HEX
        theme_str = state.get("theme") or "light"
        return cls(accent_hex=accent, is_dark=(theme_str == "dark"))


# ── Typography helpers ────────────────────────────────────────────────────────

def _set_run_font(run, size_pt: float, bold: bool = False, italic: bool = False,
                  color: Optional[RGBColor] = None, theme: Optional["_Theme"] = None):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or (theme.fg if theme else _rgb(0x0A, 0x0A, 0x0A))
    try:
        run.font.name = _FONT
    except Exception:
        try:
            run.font.name = _FONT_FALLBACK
        except Exception:
            pass


def _parse_md_segments(text: str) -> list[tuple[str, bool]]:
    """Split text on **bold** markers → [(segment, is_bold), ...]."""
    parts = re.split(r'\*\*(.+?)\*\*', text)
    return [(p, i % 2 == 1) for i, p in enumerate(parts) if p]


def _add_md_para(tf, text: str, size_pt: float, t: "_Theme",
                 color: Optional[RGBColor] = None,
                 bold_color: Optional[RGBColor] = None,
                 align: PP_ALIGN = PP_ALIGN.LEFT,
                 space_before_pt: float = 0,
                 first_para: bool = False) -> None:
    """One paragraph, bold/normal segments all inline."""
    para = tf.paragraphs[0] if first_para else tf.add_paragraph()
    para.alignment = align
    if space_before_pt:
        para.space_before = Pt(space_before_pt)
    base_color = color or t.fg
    b_color = bold_color or t.accent
    for segment, is_bold in _parse_md_segments(text):
        run = para.add_run()
        run.text = segment
        _set_run_font(run, size_pt, bold=is_bold,
                      color=b_color if is_bold else base_color, theme=t)


def _add_textbox(slide, left, top, width, height, text: str,
                 size_pt: float, t: "_Theme",
                 bold: bool = False, italic: bool = False,
                 color: Optional[RGBColor] = None,
                 align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    _set_run_font(run, size_pt, bold=bold, italic=italic, color=color or t.fg, theme=t)


# ── Structural helpers ────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _accent_bar(slide, t: "_Theme") -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.1), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = t.accent
    bar.line.fill.background()


def _title_bar(slide, title: str, t: "_Theme", top: float = Inches(0.22)) -> None:
    """Title textbox (no line overlap: line starts AFTER the text area)."""
    # Title box is tall enough for 2 lines
    _add_textbox(slide,
        left=Inches(0.9), top=top, width=Inches(11.5), height=Inches(1.0),
        text=title, size_pt=20, bold=True, color=t.fg, t=t,
    )
    # Divider line placed below the title box
    line = slide.shapes.add_shape(
        1, Inches(0.9), top + Inches(1.05), Inches(11.5), Inches(0.015),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = t.accent
    line.line.fill.background()


def _draw_veille_logo(slide, cx, cy, radius, t: "_Theme",
                      on_accent_bg: bool = False) -> None:
    """
    Draw the Veille radar logo at (cx, cy) with given radius.
    Structure mirrors Logo.tsx:
      - 3 concentric rings (stroke, no fill)
      - 3 satellite dots (chart palette colors)
      - 3 thin lines from center to each satellite (rotated rectangles — no connectors)
      - 1 center dot (accent color)
    All position/size values are EMUs.
    NOTE: We deliberately avoid add_connector() — it produces unresolvable shape-ID
    references in the XML that cause PowerPoint's repair dialog on open.
    """
    ring_color = t.white if on_accent_bg else t.muted
    line_color = t.white if on_accent_bg else t.border

    def _oval(x_c, y_c, r, fill_color: Optional[RGBColor], stroke_color: RGBColor,
              stroke_pt: float = 0.75):
        shape = slide.shapes.add_shape(9, x_c - r, y_c - r, r * 2, r * 2)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        shape.line.color.rgb = stroke_color
        shape.line.width = Pt(stroke_pt)

    def _line(x1, y1, x2, y2, color: RGBColor, thickness_pt: float = 0.4):
        """Draw a line as a rotated thin rectangle — compatible with all PPTX viewers."""
        dx = x2 - x1
        dy = y2 - y1
        length = math.hypot(dx, dy)
        if length < 1:
            return
        angle_deg = math.degrees(math.atan2(dy, dx))
        # Height of the rectangle (line thickness)
        h = max(int(Pt(thickness_pt)), 9144)  # min 1pt
        # Rectangle left/top so its left-center is at (x1, y1) before rotation
        # python-pptx rotates around the shape center, so offset accordingly
        rect_left = x1
        rect_top  = y1 - h // 2
        shape = slide.shapes.add_shape(1, rect_left, rect_top, int(length), h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.rotation = angle_deg

    # 3 concentric rings at 100%, 78%, 56% of radius
    for scale in (1.0, 0.78, 0.56):
        _oval(cx, cy, int(radius * scale), None, ring_color, stroke_pt=0.6)

    # Satellite positions (match Logo.tsx viewBox 40×40, center at 20,20)
    # Logo.tsx: (28,12), (32,24), (12,28) → offsets from center: (+8,-8),(+12,+4),(-8,+8)
    # Normalised by half-viewbox (20): 0.40, 0.60, etc.
    vp = radius
    satellites = [
        (cx + int(vp * 0.40), cy - int(vp * 0.40), t.chart_palette[4 % len(t.chart_palette)]),
        (cx + int(vp * 0.60), cy + int(vp * 0.20), t.chart_palette[1 % len(t.chart_palette)]),
        (cx - int(vp * 0.40), cy + int(vp * 0.40), t.chart_palette[3 % len(t.chart_palette)]),
    ]
    dot_r = max(int(radius * 0.10), 30000)

    # Lines from center to each satellite (draw before dots so dots sit on top)
    for sx, sy, _ in satellites:
        _line(cx, cy, sx, sy, line_color)

    # Satellite dots
    for sx, sy, color in satellites:
        _oval(sx, sy, dot_r, color, color, stroke_pt=0)

    # Center dot (accent color, drawn last so it's on top)
    center_r = max(int(radius * 0.30), 60000)
    _oval(cx, cy, center_r, t.accent, t.accent, stroke_pt=0)


def _footer(slide, t: "_Theme", on_accent_bg: bool = False) -> None:
    """Logo + 'Powered by Veille' anchored to bottom-right corner."""
    logo_r = Inches(0.115)   # radius of whole logo
    # Right-align: logo right edge = slide right - 0.18"
    logo_cx = SLIDE_W - Inches(0.18) - logo_r
    logo_cy = SLIDE_H - Inches(0.22)

    _draw_veille_logo(slide, logo_cx, logo_cy, logo_r, t, on_accent_bg=on_accent_bg)

    # Text to the left of the logo, right-aligned to slide edge
    text_color = (t.white if on_accent_bg else t.muted)
    text_w = Inches(1.85)
    _add_textbox(slide,
        left=logo_cx - logo_r - text_w - Inches(0.04),
        top=SLIDE_H - Inches(0.35),
        width=text_w, height=Inches(0.3),
        text="Powered by Veille",
        size_pt=7.5, color=text_color, t=t, align=PP_ALIGN.RIGHT,
    )


# ── Chart styling ─────────────────────────────────────────────────────────────

def _remove_chart_bg(chart) -> None:
    """Force transparent fill on chart space and plot area."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree

        def _no_fill(parent, sp_pr_tag: str):
            sp_pr = parent.find(qn(sp_pr_tag))
            if sp_pr is None:
                sp_pr = etree.SubElement(parent, qn(sp_pr_tag))
            for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:pattFill'):
                el = sp_pr.find(qn(tag))
                if el is not None:
                    sp_pr.remove(el)
            sp_pr.insert(0, etree.Element(qn('a:noFill')))

        _no_fill(chart._element, 'c:spPr')
        c_chart = chart._element.find(qn('c:chart'))
        if c_chart is not None:
            pa = c_chart.find(qn('c:plotArea'))
            if pa is not None:
                _no_fill(pa, 'c:spPr')
    except Exception as e:
        logger.debug("Chart bg removal failed: %s", e)


def _style_chart_axes(chart, t: "_Theme") -> None:
    """Muted axis labels + light gridlines — readable on any background."""
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        label_hex = 'A1A1A1' if t.is_dark else '737373'
        grid_hex  = '2E2E2E' if t.is_dark else 'E5E5E5'

        c_chart = chart._element.find(qn('c:chart'))
        if c_chart is None:
            return
        pa = c_chart.find(qn('c:plotArea'))
        if pa is None:
            return

        for ax_tag in ('c:valAx', 'c:catAx', 'c:dateAx', 'c:serAx'):
            for ax in pa.findall(qn(ax_tag)):
                # Label color
                tx_pr = ax.find(qn('c:txPr'))
                if tx_pr is None:
                    tx_pr = etree.SubElement(ax, qn('c:txPr'))
                    etree.SubElement(tx_pr, qn('a:bodyPr'))
                    etree.SubElement(tx_pr, qn('a:lstStyle'))
                p = tx_pr.find(qn('a:p'))
                if p is None:
                    p = etree.SubElement(tx_pr, qn('a:p'))
                r = p.find(qn('a:r'))
                if r is None:
                    r = etree.SubElement(p, qn('a:r'))
                rpr = r.find(qn('a:rPr'))
                if rpr is None:
                    rpr = etree.SubElement(r, qn('a:rPr'))
                rpr.set('sz', '800')
                for sf in rpr.findall(qn('a:solidFill')):
                    rpr.remove(sf)
                sf = etree.SubElement(rpr, qn('a:solidFill'))
                etree.SubElement(sf, qn('a:srgbClr')).set('val', label_hex)

                # Gridlines
                for gl in ax.findall(qn('c:majorGridlines')):
                    sp = gl.find(qn('c:spPr'))
                    if sp is None:
                        sp = etree.SubElement(gl, qn('c:spPr'))
                    ln = sp.find(qn('a:ln'))
                    if ln is None:
                        ln = etree.SubElement(sp, qn('a:ln'))
                    for sf2 in ln.findall(qn('a:solidFill')):
                        ln.remove(sf2)
                    sf2 = etree.SubElement(ln, qn('a:solidFill'))
                    etree.SubElement(sf2, qn('a:srgbClr')).set('val', grid_hex)

                # Axis line — invisible
                ax_sp = ax.find(qn('c:spPr'))
                if ax_sp is None:
                    ax_sp = etree.SubElement(ax, qn('c:spPr'))
                ax_ln = ax_sp.find(qn('a:ln'))
                if ax_ln is None:
                    ax_ln = etree.SubElement(ax_sp, qn('a:ln'))
                if ax_ln.find(qn('a:noFill')) is None and ax_ln.find(qn('a:solidFill')) is None:
                    etree.SubElement(ax_ln, qn('a:noFill'))
    except Exception as e:
        logger.debug("Axis styling failed: %s", e)


def _color_chart_series(chart, palette: list[RGBColor]) -> None:
    try:
        plot = chart.plots[0]
        # Pie/doughnut: color individual points
        if hasattr(plot, 'series') and plot.series:
            try:
                for i, point in enumerate(plot.series[0].points):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = palette[i % len(palette)]
                return
            except Exception:
                pass
        for i, series in enumerate(chart.series):
            c = palette[i % len(palette)]
            try:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = c
            except Exception:
                pass
            try:
                series.format.line.color.rgb = c
            except Exception:
                pass
    except Exception as e:
        logger.debug("Chart color failed: %s", e)


def _style_chart_legend(chart, t: "_Theme") -> None:
    try:
        chart.has_legend = True
        chart.legend.position = 2  # BOTTOM
        chart.legend.include_in_layout = False
        from pptx.oxml.ns import qn
        from lxml import etree
        legend_hex = 'A1A1A1' if t.is_dark else '737373'
        leg = chart.legend._element
        tx_pr = leg.find(qn('c:txPr'))
        if tx_pr is None:
            tx_pr = etree.SubElement(leg, qn('c:txPr'))
            etree.SubElement(tx_pr, qn('a:bodyPr'))
            etree.SubElement(tx_pr, qn('a:lstStyle'))
        p = tx_pr.find(qn('a:p'))
        if p is None:
            p = etree.SubElement(tx_pr, qn('a:p'))
        r = p.find(qn('a:r'))
        if r is None:
            r = etree.SubElement(p, qn('a:r'))
        rpr = r.find(qn('a:rPr'))
        if rpr is None:
            rpr = etree.SubElement(r, qn('a:rPr'))
        rpr.set('sz', '900')
        for sf in rpr.findall(qn('a:solidFill')):
            rpr.remove(sf)
        sf = etree.SubElement(rpr, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr')).set('val', legend_hex)
    except Exception:
        pass


def _apply_chart_style(chart, t: "_Theme", palette_offset: int = 0) -> None:
    palette = t.chart_palette[palette_offset:] + t.chart_palette[:palette_offset]
    _remove_chart_bg(chart)
    _color_chart_series(chart, palette)
    _style_chart_axes(chart, t)
    _style_chart_legend(chart, t)


def _make_chart(slide, chart_type, chart_data, left, top, width, height,
                t: "_Theme", palette_offset: int = 0):
    shape = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    _apply_chart_style(shape.chart, t, palette_offset)
    return shape


# ── Slide renderers ───────────────────────────────────────────────────────────

def _render_title_slide(slide, spec: dict, t: "_Theme") -> None:
    _set_bg(slide, t.bg)
    _accent_bar(slide, t)

    # Top accent strip
    strip = slide.shapes.add_shape(1, Inches(0.1), Inches(0), SLIDE_W - Inches(0.1), Inches(0.06))
    strip.fill.solid()
    strip.fill.fore_color.rgb = t.accent
    strip.line.fill.background()

    title = spec.get("title", "Presentation")
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(10.8), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    _add_md_para(tf, title, size_pt=38, t=t, color=t.fg, bold_color=t.accent, first_para=True)

    subtitle = spec.get("subtitle", "")
    if subtitle:
        _add_textbox(slide,
            left=Inches(1.0), top=Inches(4.3), width=Inches(9.5), height=Inches(0.75),
            text=subtitle, size_pt=15, italic=True, color=t.accent, t=t,
        )

    div = slide.shapes.add_shape(1, Inches(1.0), Inches(4.1), Inches(6.0), Inches(0.015))
    div.fill.solid()
    div.fill.fore_color.rgb = t.border
    div.line.fill.background()

    _footer(slide, t)


def _render_section(slide, spec: dict, t: "_Theme") -> None:
    """Section divider — accent-colored background."""
    _set_bg(slide, t.accent)

    title = spec.get("title", "")
    subtitle = spec.get("subtitle", "")

    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.4), Inches(11.0), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    _add_md_para(tf, title, size_pt=34, t=t, color=t.white, first_para=True)

    if subtitle:
        _add_textbox(slide,
            left=Inches(1.2), top=Inches(4.4), width=Inches(9.0), height=Inches(0.7),
            text=subtitle, size_pt=15, color=t.white, t=t,
        )

    _footer(slide, t, on_accent_bg=True)


def _render_bullets(slide, spec: dict, t: "_Theme") -> None:
    _set_bg(slide, t.bg)
    _accent_bar(slide, t)
    _title_bar(slide, spec.get("title", ""), t)

    bullets = spec.get("bullets", [])
    if not bullets:
        _footer(slide, t)
        return

    # Content area starts below title divider (title top 0.22 + box 1.0 + line 0.015 + gap)
    txBox = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(11.5), Inches(5.65))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        clean = re.sub(r'^[•\-\*]\s*', '', str(bullet).strip())
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(10) if i > 0 else Pt(0)

        marker = para.add_run()
        marker.text = "▸  "
        _set_run_font(marker, 13, bold=True, color=t.accent, theme=t)

        for segment, is_bold in _parse_md_segments(clean):
            run = para.add_run()
            run.text = segment
            _set_run_font(run, 14, bold=is_bold,
                          color=t.accent_dark if is_bold else t.fg, theme=t)

    _footer(slide, t)


def _render_kpi_grid(slide, spec: dict, t: "_Theme") -> None:
    _set_bg(slide, t.bg)
    _accent_bar(slide, t)
    _title_bar(slide, spec.get("title", "Key Metrics"), t)

    items = spec.get("items", [])
    if not items:
        _footer(slide, t)
        return

    items = items[:8]
    count = len(items)
    per_row = min(count, 4)
    rows = (count + per_row - 1) // per_row

    card_h = Inches(2.2) if rows == 1 else Inches(2.0)
    card_w = Inches(10.8 / per_row)
    start_x = Inches(0.9)
    start_y = Inches(1.55)
    v_gap = Inches(0.15)

    for idx, item in enumerate(items):
        row = idx // per_row
        col = idx % per_row
        x = start_x + card_w * col + Inches(0.05)
        y = start_y + (card_h + v_gap) * row

        card = slide.shapes.add_shape(1, x, y, card_w - Inches(0.1), card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = t.card
        card.line.color.rgb = t.border
        card.line.width = Pt(0.75)

        strip = slide.shapes.add_shape(1, x, y, card_w - Inches(0.1), Inches(0.04))
        strip.fill.solid()
        strip.fill.fore_color.rgb = t.accent
        strip.line.fill.background()

        _add_textbox(slide,
            left=x + Inches(0.12), top=y + Inches(0.2),
            width=card_w - Inches(0.34), height=Inches(1.1),
            text=str(item.get("value", "")), size_pt=28, bold=True, color=t.accent, t=t,
            align=PP_ALIGN.CENTER,
        )
        _add_textbox(slide,
            left=x + Inches(0.1), top=y + Inches(1.38),
            width=card_w - Inches(0.3), height=Inches(0.6),
            text=str(item.get("label", "")), size_pt=11, color=t.muted, t=t,
            align=PP_ALIGN.CENTER,
        )

    _footer(slide, t)


def _render_chart_bar(slide, spec: dict, t: "_Theme") -> None:
    _set_bg(slide, t.bg)
    _accent_bar(slide, t)
    _title_bar(slide, spec.get("title", ""), t)

    labels = [str(l) for l in spec.get("labels", [])]
    values = [float(v) if v is not None else 0 for v in spec.get("values", [])]
    if labels and values and len(labels) == len(values):
        cd = CategoryChartData()
        cd.categories = labels
        cd.add_series("", values)
        _make_chart(slide, XL_CHART_TYPE.BAR_CLUSTERED, cd,
                    Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.5), t)
    _footer(slide, t)


def _render_chart_pie(slide, spec: dict, t: "_Theme") -> None:
    _set_bg(slide, t.bg)
    _accent_bar(slide, t)
    _title_bar(slide, spec.get("title", ""), t)

    labels = [str(l) for l in spec.get("labels", [])]
    values = [float(v) if v is not None else 0 for v in spec.get("values", [])]
    if labels and values and len(labels) == len(values):
        cd = CategoryChartData()
        cd.categories = labels
        cd.add_series("", values)
        # Leave generous height so legend fits below pie without overlap
        _make_chart(slide, XL_CHART_TYPE.PIE, cd,
                    Inches(2.0), Inches(1.5), Inches(9.0), Inches(5.5), t)
    _footer(slide, t)


def _render_chart_line(slide, spec: dict, t: "_Theme") -> None:
    _set_bg(slide, t.bg)
    _accent_bar(slide, t)
    _title_bar(slide, spec.get("title", ""), t)

    dates = [str(d) for d in spec.get("dates", [])]
    values = [float(v) if v is not None else 0 for v in spec.get("values", [])]
    if dates and values and len(dates) == len(values):
        cd = CategoryChartData()
        cd.categories = dates
        cd.add_series("Volume", values)
        _make_chart(slide, XL_CHART_TYPE.LINE, cd,
                    Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.5), t)
    _footer(slide, t)


def _render_chart_row(slide, spec: dict, t: "_Theme") -> None:
    """Two charts side-by-side."""
    _set_bg(slide, t.bg)
    _accent_bar(slide, t)

    title = spec.get("title", "")
    if title:
        _title_bar(slide, title, t)
        chart_top = Inches(1.5)
    else:
        chart_top = Inches(0.35)

    charts = spec.get("charts", [])[:2]
    chart_w = Inches(5.7)
    chart_h = SLIDE_H - chart_top - Inches(0.55)
    positions = [(Inches(0.85), chart_top), (Inches(7.0), chart_top)]

    for i, cs in enumerate(charts):
        lx, ty = positions[i]
        sub = cs.get("title", "")
        if sub:
            _add_textbox(slide, left=lx, top=ty - Inches(0.30), width=chart_w,
                         height=Inches(0.28), text=sub, size_pt=11, bold=True,
                         color=t.muted, t=t)

        labels = [str(l) for l in cs.get("labels", [])]
        values = [float(v) if v is not None else 0 for v in cs.get("values", [])]
        dates  = [str(d) for d in cs.get("dates", [])]
        cats   = labels or dates

        if not cats or not values:
            continue

        cd = CategoryChartData()
        cd.categories = cats
        cd.add_series("", values)

        xl_type = {
            "chart_bar": XL_CHART_TYPE.BAR_CLUSTERED, "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "chart_pie": XL_CHART_TYPE.PIE,            "pie": XL_CHART_TYPE.PIE,
            "chart_line": XL_CHART_TYPE.LINE,          "line": XL_CHART_TYPE.LINE,
        }.get(cs.get("type", ""), XL_CHART_TYPE.BAR_CLUSTERED)

        _make_chart(slide, xl_type, cd, lx, ty, chart_w, chart_h, t, palette_offset=i * 2)

    _footer(slide, t)


def _render_table(slide, spec: dict, t: "_Theme") -> None:
    _set_bg(slide, t.bg)
    _accent_bar(slide, t)
    _title_bar(slide, spec.get("title", ""), t)

    columns = spec.get("columns", [])
    rows    = spec.get("rows", [])
    if not columns or not rows:
        _footer(slide, t)
        return

    num_cols = len(columns)
    num_rows = min(len(rows), 12) + 1

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.9), Inches(1.5),
        Inches(11.5), Inches(min(5.5, num_rows * 0.43)),
    )
    tbl = tbl_shape.table

    for ci, col_name in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = t.accent
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.text = str(col_name)
        _set_run_font(run, 11, bold=True, color=t.white, theme=t)

    for ri, row_data in enumerate(rows[:11]):
        bg = t.card if ri % 2 == 0 else t.surface
        for ci in range(num_cols):
            cell = tbl.cell(ri + 1, ci)
            val = row_data[ci] if ci < len(row_data) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            para = cell.text_frame.paragraphs[0]
            for run in para.runs:
                run.text = ""
            run = para.runs[0] if para.runs else para.add_run()
            run.text = str(val) if val is not None else ""
            _set_run_font(run, 10, color=t.fg, theme=t)

    _footer(slide, t)


def _render_key_finding(slide, spec: dict, t: "_Theme") -> None:
    _set_bg(slide, t.bg)
    _accent_bar(slide, t)
    _title_bar(slide, spec.get("title", "Key Finding"), t)

    significance = spec.get("significance", "notable")
    finding = spec.get("finding", "")

    # Card
    card = slide.shapes.add_shape(1, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = t.card
    card.line.color.rgb = t.accent if significance == "surprising" else t.border
    card.line.width = Pt(1.0 if significance == "surprising" else 0.75)

    # Accent strip atop card
    strip = slide.shapes.add_shape(1, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.045))
    strip.fill.solid()
    strip.fill.fore_color.rgb = t.accent
    strip.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1.15), Inches(1.72), Inches(11.0), Inches(4.7))
    tf = txBox.text_frame
    tf.word_wrap = True

    badge = "⚡ SURPRISING SIGNAL" if significance == "surprising" else "📌 KEY FINDING"
    para0 = tf.paragraphs[0]
    br = para0.add_run()
    br.text = badge
    _set_run_font(br, 9, bold=True, color=t.accent, theme=t)

    # Finding body — ALL segments inline on one paragraph
    if finding:
        para = tf.add_paragraph()
        para.space_before = Pt(14)
        for segment, is_bold in _parse_md_segments(finding):
            run = para.add_run()
            run.text = segment
            _set_run_font(run, 17, bold=is_bold,
                          color=t.accent_dark if is_bold else t.fg, theme=t)

    _footer(slide, t)


def _render_closing(slide, spec: dict, t: "_Theme") -> None:
    _set_bg(slide, t.bg)

    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.07))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = t.accent
    top_bar.line.fill.background()

    _accent_bar(slide, t)

    title = spec.get("title", "Thank you")
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(11.0), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    _add_md_para(tf, title, size_pt=38, t=t, color=t.fg,
                 align=PP_ALIGN.CENTER, first_para=True)

    message = spec.get("message", "")
    if message:
        _add_textbox(slide,
            left=Inches(1.5), top=Inches(3.95), width=Inches(10.5), height=Inches(1.0),
            text=message, size_pt=15, color=t.accent, t=t, align=PP_ALIGN.CENTER,
        )

    # Centered logo + wordmark
    logo_r = Inches(0.14)
    logo_cx = SLIDE_W / 2 - Inches(0.9)
    logo_cy = SLIDE_H - Inches(0.65)
    _draw_veille_logo(slide, logo_cx, logo_cy, logo_r, t)
    _add_textbox(slide,
        left=logo_cx + logo_r + Inches(0.06), top=SLIDE_H - Inches(0.8),
        width=Inches(1.85), height=Inches(0.45),
        text="Powered by Veille", size_pt=10, color=t.muted, t=t,
    )


# ── Dispatch map ──────────────────────────────────────────────────────────────

_RENDERERS = {
    "title_slide": _render_title_slide,
    "section":     _render_section,
    "bullets":     _render_bullets,
    "kpi_grid":    _render_kpi_grid,
    "chart_bar":   _render_chart_bar,
    "chart_pie":   _render_chart_pie,
    "chart_line":  _render_chart_line,
    "chart_row":   _render_chart_row,
    "table":       _render_table,
    "key_finding": _render_key_finding,
    "closing":     _render_closing,
}


# ── GCS helpers ───────────────────────────────────────────────────────────────

def _get_bucket_name() -> str:
    settings = get_settings()
    return settings.gcs_presentations_bucket or settings.gcs_exports_bucket


def _upload_to_gcs(blob_name: str, data: bytes) -> str:
    client = get_gcs()
    bucket = client.bucket(_get_bucket_name())
    blob = bucket.blob(blob_name)
    blob.upload_from_string(
        data,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    return f"gs://{_get_bucket_name()}/{blob_name}"


def _download_template(gcs_path: str) -> bytes:
    client = get_gcs()
    if gcs_path.startswith("gs://"):
        without_scheme = gcs_path[5:]
        bucket_name, _, blob_name = without_scheme.partition("/")
    else:
        bucket_name = _get_bucket_name()
        blob_name = gcs_path
    return client.bucket(bucket_name).blob(blob_name).download_as_bytes()


def _find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.lower() in ('blank', 'leer', 'vide', 'vuoto'):
            return layout
    for layout in prs.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout
    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


# ── Main tool ─────────────────────────────────────────────────────────────────

def generate_presentation(
    collection_ids: list[str],
    slides: list[dict],
    title: str = "",
    template_gcs_path: str = "",
    tool_context=None,
) -> dict:
    """Render a PowerPoint presentation from a structured slide spec.

    WHEN TO USE: When the user wants a downloadable presentation deck based on
    collected data. Always gather data first (get_collection_stats, execute_sql),
    then design the slides and call this tool.

    CONTEXT-ADAPTIVE DESIGN — think like a senior analyst deciding what the
    audience needs:
    - Read the session: what data was collected, what queries were run, what
      the user's actual question was. Let THAT determine the slide structure.
    - A simple question deserves 3-4 focused slides. A deep comparative study
      may warrant 8-10. Never pad — never cut insights to hit a target count.
    - Only include slide types that add new information not already visible in
      other slides. Don't repeat the same insight in different formats.
    - VARY the structure based on what the data shows:
      * Sentiment-heavy data → chart_pie + key_finding + bullets
      * Time-series data → chart_line + key_finding
      * Comparative data → chart_bar + chart_row + table
      * Few clear KPIs → kpi_grid + closing
      * Complex narrative → bullets + key_finding + closing
    - If the data is surprising, use key_finding with significance "surprising".
      If the data is ordinary, skip key_finding entirely.

    SLIDE TYPES SUPPORTED:
    - title_slide:  {type, title, subtitle}
    - section:      {type, title, subtitle?}  — visual section divider (use only
                    for decks with 7+ slides, to break into logical chapters)
    - bullets:      {type, title, bullets: [str, ...]}  — supports **bold** markdown
    - kpi_grid:     {type, title, items: [{label, value}, ...]}  — up to 8 KPIs
    - chart_bar:    {type, title, labels: [str], values: [num]}
    - chart_pie:    {type, title, labels: [str], values: [num]}
    - chart_line:   {type, title, dates: [str], values: [num]}
    - chart_row:    {type, title?, charts: [{type, title, labels/dates, values}]}
                    — 2 related charts side-by-side (saves a slide)
    - table:        {type, title, columns: [str], rows: [[val, ...], ...]}
    - key_finding:  {type, title, finding: str, significance: "surprising"|"notable"}
    - closing:      {type, title, message?}

    MARKDOWN: Use **bold** in bullet text and key_finding finding text.
    CHART_ROW: Prefer chart_row when showing 2 related distributions (e.g.
    sentiment split + platform split) — one slide, two charts.

    Args:
        collection_ids: Collections this presentation covers.
        slides: Ordered list of slide specs (see types above).
        title: Presentation title.
        template_gcs_path: Optional GCS path to a user-supplied .pptx template.

    Returns:
        status, presentation_id, title, collection_ids, slide_count on success.
    """
    if not slides:
        return {"status": "error", "message": "slides list is empty."}

    # Derive theme from session state
    try:
        from google.adk.tools.tool_context import ToolContext as _TC  # noqa
        _ctx = tool_context
    except ImportError:
        _ctx = None
    theme = _Theme.from_session(_ctx)

    # ── Load base presentation ────────────────────────────────────────────────
    try:
        if template_gcs_path:
            template_bytes = _download_template(template_gcs_path)
            prs = Presentation(io.BytesIO(template_bytes))
        else:
            template_path = (
                Path(__file__).parent.parent.parent / "api" / "assets" / "templates" / "default.pptx"
            )
            prs = Presentation(str(template_path))

        # Remove template slides properly: drop both the package part (so the
        # XML file is gone from the ZIP) AND the sldId reference.  Just removing
        # from _sldIdLst leaves the slide XML files in the package, causing
        # duplicate ZIP entries that make PowerPoint show its repair dialog.
        xml_slides = prs.slides._sldIdLst
        rIds = [el.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                for el in list(xml_slides)]
        for rId in rIds:
            if rId:
                try:
                    prs.part.drop_rel(rId)
                except Exception:
                    pass
        for sld_id in list(xml_slides):
            xml_slides.remove(sld_id)
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
    except Exception as e:
        logger.warning("generate_presentation: failed to load template: %s", e)
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    blank_layout = _find_blank_layout(prs)

    # ── Render slides ─────────────────────────────────────────────────────────
    rendered = 0
    for spec in slides:
        slide_type = spec.get("type", "")
        renderer = _RENDERERS.get(slide_type)
        if renderer is None:
            logger.warning("generate_presentation: unknown slide type %r — skipping", slide_type)
            continue
        try:
            slide = prs.slides.add_slide(blank_layout)
            renderer(slide, spec, theme)
            rendered += 1
        except Exception as e:
            logger.warning("generate_presentation: failed to render %r: %s", slide_type, e)

    if rendered == 0:
        return {"status": "error", "message": "No slides could be rendered — check the slides spec."}

    # ── Serialize and upload ──────────────────────────────────────────────────
    presentation_id = f"ppt-{uuid.uuid4().hex[:10]}"
    safe_title = (title or "presentation").replace(" ", "_").replace("/", "-")[:60]
    blob_name = f"presentations/{presentation_id}/{safe_title}.pptx"

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    try:
        _upload_to_gcs(blob_name, pptx_bytes)
        logger.info("generate_presentation: uploaded %s (%d bytes)", blob_name, len(pptx_bytes))
    except Exception as e:
        logger.error("generate_presentation: GCS upload failed: %s", e)
        return {"status": "error", "message": f"Failed to save presentation: {e}"}

    return {
        "status": "success",
        "presentation_id": presentation_id,
        "title": title or "Presentation",
        "collection_ids": collection_ids,
        "slide_count": rendered,
        "gcs_path": blob_name,
        "message": f"Presentation created with {rendered} slides. The user can download it from the Artifacts panel.",
    }
