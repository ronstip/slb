"""Layout matching — resolve requested layout names to template layout indices."""

import logging
from typing import Any, Optional

logger = logging.getLogger(__name__)

# Normalized aliases for localized layout names
# Includes Google Slides layout names (SCREAMING_SNAKE_CASE)
_LAYOUT_ALIASES: dict[str, list[str]] = {
    "title slide": [
        "titelfolie", "diapositive de titre", "diapositiva del titulo",
        "title", "cover", "opening",
    ],
    "title and content": [
        "titel und inhalt", "titre et contenu", "titulo y contenido",
        "content", "title + content",
        "title_and_body", "one_column_text",  # Google Slides
    ],
    "section header": [
        "abschnittsheader", "en-tete de section", "encabezado de seccion",
        "section", "divider", "section divider",
        "section_header",  # Google Slides
    ],
    "two content": [
        "zwei inhalte", "deux contenus", "dos contenidos",
        "two column", "side by side", "split",
        "title_and_two_columns",  # Google Slides
    ],
    "comparison": [
        "vergleich", "comparaison", "comparacion",
        "compare",
    ],
    "title only": [
        "nur titel", "titre seul", "solo titulo",
        "title_only",  # Google Slides
    ],
    "blank": [
        "leer", "vide", "vuoto", "en blanco",
    ],
    "content with caption": [
        "inhalt mit beschriftung",
        "caption_only",  # Google Slides
    ],
    "picture with caption": [
        "bild mit beschriftung",
    ],
    "big number": [
        "big_number",  # Google Slides
    ],
    "main point": [
        "main_point",  # Google Slides
    ],
}


def _normalize(name: str) -> str:
    return name.strip().lower()


def _build_signature(placeholders: list[dict]) -> frozenset:
    """Build a structural signature from placeholder types.

    Returns a frozenset of (type_name, count) tuples for content placeholders.
    """
    type_counts: dict[str, int] = {}
    for ph in placeholders:
        ph_type = ph.get("type", "UNKNOWN")
        if ph_type in ("DATE", "FOOTER", "SLIDE_NUMBER", "HEADER", "UNKNOWN"):
            continue
        # Normalize CENTER_TITLE to TITLE for signature matching
        if ph_type == "CENTER_TITLE":
            ph_type = "TITLE"
        type_counts[ph_type] = type_counts.get(ph_type, 0) + 1
    return frozenset(type_counts.items())


# Known structural signatures -> canonical layout names
_KNOWN_SIGNATURES = {
    frozenset([("TITLE", 1), ("SUBTITLE", 1)]): "Title Slide",
    frozenset([("TITLE", 1), ("OBJECT", 1)]): "Title and Content",
    frozenset([("TITLE", 1), ("BODY", 1)]): "Title and Content",  # Google Slides uses BODY
    frozenset([("TITLE", 1), ("OBJECT", 2)]): "Two Content",
    frozenset([("TITLE", 1), ("BODY", 2)]): "Two Content",  # Google Slides uses BODY
    frozenset([("TITLE", 1), ("BODY", 2), ("OBJECT", 2)]): "Comparison",
    frozenset([("TITLE", 1)]): "Title Only",
    frozenset(): "Blank",
}


def resolve_layout(
    requested_name: str,
    manifest_layouts: list[dict],
) -> Optional[int]:
    """Resolve a requested layout name to a template layout index.

    Matching priority:
    1. Exact name match (case-insensitive)
    2. Alias match
    3. Structural signature match
    4. None (caller should fall back to Blank)

    Returns:
        Layout index or None if no match found.
    """
    req_norm = _normalize(requested_name)

    # 1. Exact name match
    for layout in manifest_layouts:
        if _normalize(layout["name"]) == req_norm:
            return layout["index"]

    # 2. Alias match — find the alias group, then check if the canonical name
    # OR any alias in that group matches a template layout name.
    matched_group = None
    for canon, aliases in _LAYOUT_ALIASES.items():
        if req_norm == canon or req_norm in aliases:
            matched_group = (canon, aliases)
            break

    if matched_group:
        canon, aliases = matched_group
        all_names = {canon} | set(aliases)
        for layout in manifest_layouts:
            if _normalize(layout["name"]) in all_names:
                return layout["index"]

    # 3. Structural signature match
    # First determine what signature the requested name expects
    matched_canon = matched_group[0] if matched_group else None
    target_sig = None
    for sig, canon in _KNOWN_SIGNATURES.items():
        if _normalize(canon) == req_norm or (matched_canon and _normalize(canon) == matched_canon):
            target_sig = sig
            break

    if target_sig is not None:
        for layout in manifest_layouts:
            layout_sig = _build_signature(layout.get("placeholders", []))
            if layout_sig == target_sig:
                return layout["index"]

    # 4. No match
    logger.warning("resolve_layout: no match for %r in template", requested_name)
    return None


def find_blank_layout(manifest_layouts: list[dict]) -> int:
    """Find the Blank layout index, or the last layout as fallback."""
    for layout in manifest_layouts:
        if _normalize(layout["name"]) in ("blank", "leer", "vide", "vuoto", "en blanco"):
            return layout["index"]
    # Fallback: return last layout index
    if manifest_layouts:
        return manifest_layouts[-1]["index"]
    return 0


def get_slot_placeholders(layout_info: dict) -> dict[str, dict]:
    """Map slot names to their placeholder info for a given layout.

    Returns: {"title": {idx, type, left, top, width, height}, "body": {...}, ...}
    """
    result = {}
    for ph in layout_info.get("placeholders", []):
        slot = ph.get("slot", "")
        if slot:
            result[slot] = ph
    return result


def compute_free_area(
    layout_info: dict,
    slide_width: int,
    slide_height: int,
) -> tuple[int, int, int, int]:
    """Compute the free area below the title placeholder.

    For 'Title Only' and 'Blank' layouts, this is the space where
    custom components (kpi_grid, key_finding) are rendered.

    Returns: (left, top, width, height) in EMU.
    """
    from pptx.util import Inches

    title_ph = None
    for ph in layout_info.get("placeholders", []):
        if ph.get("type") in ("TITLE", "CENTER_TITLE"):
            title_ph = ph
            break

    if title_ph:
        margin = Inches(0.35)
        top = title_ph["top"] + title_ph["height"] + margin
        left = title_ph["left"]
        width = title_ph["width"]
        height = slide_height - top - Inches(0.5)
    else:
        # No title — use most of the slide
        left = Inches(0.5)
        top = Inches(0.5)
        width = slide_width - Inches(1.0)
        height = slide_height - Inches(1.0)

    return (left, top, width, max(height, Inches(1.0)))
